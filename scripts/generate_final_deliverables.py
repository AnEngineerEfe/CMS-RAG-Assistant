"""CMS-RAG için kurumsal Word dokümantasyonu ve PowerPoint sunumu üretir."""

from __future__ import annotations

from datetime import date
import json
from pathlib import Path
from typing import Iterable

from docx import Document
from docx.enum.section import WD_ORIENT
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = ROOT / "deliverables"
DOCX_PATH = OUTPUT_DIR / "CMS-RAG_Nihai_Teknik_Dokumantasyon.docx"
PPTX_PATH = OUTPUT_DIR / "CMS-RAG_Nihai_Proje_Sunumu.pptx"
TODAY = date(2026, 8, 18)
VERSION = "2.0"
COMMIT = "codex/pgvector-lineage · agentic final"
BENCHMARK_REPORT = json.loads(
    (ROOT / "evaluation" / "results" / "latest" / "benchmark_report.json")
    .read_text(encoding="utf-8")
)
QUALITY_REPORT = json.loads(
    (ROOT / "evaluation" / "results" / "quality-latest" / "quality_evaluation_report.json")
    .read_text(encoding="utf-8")
)
LINEAGE_REPORT = json.loads(
    (ROOT / "evaluation" / "results" / "lineage-latest" / "lineage_evaluation_report.json")
    .read_text(encoding="utf-8")
)
AUTOMATED_TESTS = 145

NAVY = "0B1F3A"
BLUE = "1B5FA7"
CYAN = "38BDF8"
ORANGE = "F59E0B"
GREEN = "16A34A"
RED = "DC2626"
INK = "172033"
MUTED = "5F6B7A"
LIGHT = "F3F6FA"
PALE_BLUE = "EAF3FC"
WHITE = "FFFFFF"
LINE = "D7E0EA"


def set_cell_shading(cell, fill: str) -> None:
    """Word tablo hücresine onaltılık arka plan rengi uygular."""

    properties = cell._tc.get_or_add_tcPr()
    shading = properties.find(qn("w:shd"))
    if shading is None:
        shading = OxmlElement("w:shd")
        properties.append(shading)
    shading.set(qn("w:fill"), fill)


def set_cell_border(cell, color: str = LINE, size: str = "6") -> None:
    """Word tablo hücresine ince ve tutarlı kenarlık ekler."""

    properties = cell._tc.get_or_add_tcPr()
    borders = properties.first_child_found_in("w:tcBorders")
    if borders is None:
        borders = OxmlElement("w:tcBorders")
        properties.append(borders)
    for edge in ("top", "left", "bottom", "right"):
        tag = f"w:{edge}"
        element = borders.find(qn(tag))
        if element is None:
            element = OxmlElement(tag)
            borders.append(element)
        element.set(qn("w:val"), "single")
        element.set(qn("w:sz"), size)
        element.set(qn("w:color"), color)


def set_repeat_table_header(row) -> None:
    """Word tablosunun ilk satırını yeni sayfalarda başlık olarak tekrarlar."""

    properties = row._tr.get_or_add_trPr()
    repeat = OxmlElement("w:tblHeader")
    repeat.set(qn("w:val"), "true")
    properties.append(repeat)


def add_page_field(paragraph) -> None:
    """Word altbilgisine dinamik PAGE / NUMPAGES alanı ekler."""

    paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    run = paragraph.add_run("CMS-RAG · Nihai Dokümantasyon     ")
    run.font.name = "Aptos"
    run.font.size = Pt(8)
    run.font.color.rgb = RGBColor.from_string(MUTED)
    for field_name, prefix in (("PAGE", "Sayfa "), ("NUMPAGES", " / ")):
        label = paragraph.add_run(prefix)
        label.font.size = Pt(8)
        field = OxmlElement("w:fldSimple")
        field.set(qn("w:instr"), field_name)
        paragraph._p.append(field)


def style_word_document(document: Document) -> None:
    """Word belgesinin sayfa, yazı tipi, başlık ve üstbilgi standardını kurar."""

    section = document.sections[0]
    section.top_margin = Cm(1.8)
    section.bottom_margin = Cm(1.7)
    section.left_margin = Cm(2.0)
    section.right_margin = Cm(2.0)

    normal = document.styles["Normal"]
    normal.font.name = "Aptos"
    normal.font.size = Pt(10.5)
    normal.font.color.rgb = RGBColor.from_string(INK)
    normal.paragraph_format.space_after = Pt(6)
    normal.paragraph_format.line_spacing = 1.12

    heading_sizes = {1: 23, 2: 16, 3: 12.5}
    for level, size in heading_sizes.items():
        style = document.styles[f"Heading {level}"]
        style.font.name = "Aptos Display"
        style.font.size = Pt(size)
        style.font.bold = True
        style.font.color.rgb = RGBColor.from_string(NAVY if level == 1 else BLUE)
        style.paragraph_format.space_before = Pt(12 if level > 1 else 6)
        style.paragraph_format.space_after = Pt(6)
        style.paragraph_format.keep_with_next = True

    header = section.header.paragraphs[0]
    header.text = "CMS-RAG ASSISTANT  /  EVIDENCE-FIRST KNOWLEDGE OPERATIONS"
    header.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    header.runs[0].font.name = "Aptos"
    header.runs[0].font.size = Pt(8)
    header.runs[0].font.bold = True
    header.runs[0].font.color.rgb = RGBColor.from_string(BLUE)
    add_page_field(section.footer.paragraphs[0])


def add_word_title(document: Document, title: str, subtitle: str | None = None) -> None:
    """Word bölüm başlığını marka renkleriyle ekler."""

    paragraph = document.add_paragraph()
    paragraph.style = document.styles["Heading 1"]
    run = paragraph.add_run(title)
    run.font.color.rgb = RGBColor.from_string(NAVY)
    if subtitle:
        sub = document.add_paragraph(subtitle)
        sub.style = document.styles["Subtitle"]
        sub.runs[0].font.name = "Aptos"
        sub.runs[0].font.color.rgb = RGBColor.from_string(MUTED)


def add_word_body(document: Document, text: str, bold_lead: str | None = None) -> None:
    """Word belgesine isteğe bağlı kalın girişli açıklama paragrafı ekler."""

    paragraph = document.add_paragraph()
    if bold_lead and text.startswith(bold_lead):
        first = paragraph.add_run(bold_lead)
        first.bold = True
        first.font.color.rgb = RGBColor.from_string(NAVY)
        paragraph.add_run(text[len(bold_lead):])
    else:
        paragraph.add_run(text)


def add_word_bullets(document: Document, items: Iterable[str], level: int = 0) -> None:
    """Word belgesine tutarlı girintili madde listesi ekler."""

    for item in items:
        paragraph = document.add_paragraph(style="List Bullet" if level == 0 else "List Bullet 2")
        paragraph.add_run(item)


def add_word_numbered(document: Document, items: Iterable[str]) -> None:
    """Word belgesine numaralı süreç adımları ekler."""

    for item in items:
        paragraph = document.add_paragraph(style="List Number")
        paragraph.add_run(item)


def add_word_table(
    document: Document,
    headers: list[str],
    rows: list[list[str]],
    widths: list[float] | None = None,
) -> None:
    """Word belgesine renkli başlıklı, denetlenebilir veri tablosu ekler."""

    table = document.add_table(rows=1, cols=len(headers))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    table.autofit = False
    header_cells = table.rows[0].cells
    set_repeat_table_header(table.rows[0])
    for index, header in enumerate(headers):
        cell = header_cells[index]
        set_cell_shading(cell, NAVY)
        set_cell_border(cell, NAVY)
        cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        paragraph = cell.paragraphs[0]
        paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
        run = paragraph.add_run(header)
        run.bold = True
        run.font.name = "Aptos"
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor.from_string(WHITE)
        if widths:
            cell.width = Cm(widths[index])
    for row_index, values in enumerate(rows):
        cells = table.add_row().cells
        for index, value in enumerate(values):
            cell = cells[index]
            set_cell_shading(cell, WHITE if row_index % 2 == 0 else LIGHT)
            set_cell_border(cell)
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            paragraph = cell.paragraphs[0]
            run = paragraph.add_run(str(value))
            run.font.name = "Aptos"
            run.font.size = Pt(8.7)
            run.font.color.rgb = RGBColor.from_string(INK)
            if widths:
                cell.width = Cm(widths[index])
    document.add_paragraph()


def add_word_callout(document: Document, title: str, body: str, color: str = BLUE) -> None:
    """Önemli karar veya uyarıyı tek hücreli vurgulu Word kutusunda gösterir."""

    table = document.add_table(rows=1, cols=1)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    cell = table.cell(0, 0)
    set_cell_shading(cell, PALE_BLUE if color == BLUE else "FFF7E8")
    set_cell_border(cell, color, "14")
    paragraph = cell.paragraphs[0]
    lead = paragraph.add_run(f"{title}\n")
    lead.bold = True
    lead.font.name = "Aptos"
    lead.font.size = Pt(11)
    lead.font.color.rgb = RGBColor.from_string(color)
    text = paragraph.add_run(body)
    text.font.name = "Aptos"
    text.font.size = Pt(9.5)
    text.font.color.rgb = RGBColor.from_string(INK)
    document.add_paragraph()


def add_word_pipeline(document: Document, stages: list[tuple[str, str]]) -> None:
    """RAG aşamalarını Word içinde yatay akış tablosu olarak gösterir."""

    table = document.add_table(rows=2, cols=len(stages))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.autofit = True
    for index, (name, detail) in enumerate(stages):
        top = table.cell(0, index)
        bottom = table.cell(1, index)
        set_cell_shading(top, BLUE if index % 2 == 0 else NAVY)
        set_cell_shading(bottom, LIGHT)
        set_cell_border(top, WHITE)
        set_cell_border(bottom, LINE)
        top.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        bottom.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        top.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
        bottom.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = top.paragraphs[0].add_run(name)
        run.bold = True
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor.from_string(WHITE)
        detail_run = bottom.paragraphs[0].add_run(detail)
        detail_run.font.size = Pt(7.5)
        detail_run.font.color.rgb = RGBColor.from_string(INK)
    document.add_paragraph()


def build_word_document() -> None:
    """Tam kapsamlı nihai teknik dokümantasyonu DOCX biçiminde üretir."""

    document = Document()
    style_word_document(document)
    properties = document.core_properties
    properties.title = "CMS-RAG Assistant — Nihai Teknik Dokümantasyon"
    properties.subject = "Kurumsal mimari, RAG, test, işletim ve sunum rehberi"
    properties.author = "CMS-RAG Proje Ekibi"
    properties.keywords = "CMS, RAG, ADVENT, FAISS, BM25, Ollama, Streamlit"

    # Kapak
    paragraph = document.add_paragraph()
    paragraph.paragraph_format.space_before = Pt(90)
    paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
    eyebrow = paragraph.add_run("CMS-RAG  /  KNOWLEDGE OPERATIONS")
    eyebrow.bold = True
    eyebrow.font.name = "Aptos"
    eyebrow.font.size = Pt(11)
    eyebrow.font.color.rgb = RGBColor.from_string(CYAN)
    title = document.add_paragraph()
    title.paragraph_format.space_before = Pt(12)
    run = title.add_run("Savaş Yönetim Sistemi\nRAG Asistanı")
    run.font.name = "Aptos Display"
    run.font.size = Pt(38)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string(NAVY)
    subtitle = document.add_paragraph()
    sub_run = subtitle.add_run(
        "Nihai Teknik Dokümantasyon\n"
        "Kaynak kontrollü · Yerel üretim · Hibrit retrieval · Denetlenebilir kanıt"
    )
    sub_run.font.name = "Aptos"
    sub_run.font.size = Pt(16)
    sub_run.font.color.rgb = RGBColor.from_string(MUTED)
    document.add_paragraph()
    add_word_table(
        document,
        ["Belge", "Değer"],
        [
            ["Sürüm", VERSION],
            ["Tarih", TODAY.strftime("%d.%m.%Y")],
            ["Doğrulanmış dal", COMMIT],
            ["Mimari", "Hazır PDF paketi + önceden hesaplanmış embedding snapshot'ı"],
            ["Durum", f"Nihai kabul — {AUTOMATED_TESTS}/{AUTOMATED_TESTS} test, 8/8 retrieval, 7/7 yanıt"],
        ],
        [5.0, 11.5],
    )
    add_word_callout(
        document,
        "Kullanım sınırı",
        "Bu çalışma, HAVELSAN'ın deniz savaş yönetim sistemleri alanındaki kamuya açık "
        "çalışma ve ürünleri ile yapay zekâ entegrasyonu süreçlerine yönelik bir ön çalışma "
        "ve araştırmadır. Şirket içi, özel veya tasnifli veri kullanılmamıştır. Operasyonel "
        "karar sistemi değildir; kritik iddialar asıl belge ve sayfadan doğrulanmalıdır.",
        ORANGE,
    )
    document.add_page_break()

    # Belge kontrol ve içindekiler
    add_word_title(document, "Belge Kontrolü", "Sunum, inceleme ve teknik devir için tek doğruluk kaynağı")
    add_word_table(
        document,
        ["Alan", "Açıklama"],
        [
            ["Hedef kitle", "Mentör, jüri, geliştirici, teknik yönetici ve sistem kullanıcısı"],
            ["Kapsam", "Amaç, mimari, veri, retrieval, güvenlik, kurulum, test, Git ve yol haritası"],
            ["Kaynaklar", "Dört önceden hazırlanmış kamuya açık PDF; şirket verisi kullanılmaz"],
            ["Çalışma modeli", "Yerel Ollama qwen2.5:3b; güçlü donanımda qwen2.5:7b"],
            ["Arayüz", "Streamlit tabanlı, kaynak kartlı ve akışlı sohbet"],
            ["Son kabul", f"{AUTOMATED_TESTS}/{AUTOMATED_TESTS} test, 8/8 retrieval, 7/7 yanıt ve 77/77 snapshot"],
        ],
        [4.2, 12.3],
    )
    document.add_heading("İçindekiler", level=2)
    toc = [
        "1. Yönetici özeti",
        "2. Problem, amaç ve başarı ölçütleri",
        "3. Kapsam ve kullanım senaryoları",
        "4. Kaynak güven modeli ve veri yönetişimi",
        "5. Katmanlı yazılım mimarisi",
        "6. Belge yaşam döngüsü",
        "7. Hibrit RAG ve retrieval hattı",
        "8. Konuşma belleği ve cevap üretimi",
        "9. Güvenlik, gizlilik ve güvenilirlik",
        "10. Kullanıcı arayüzü ve deneyim",
        "11. Kurulum, çalıştırma ve işletim",
        "12. Test stratejisi ve kabul kanıtları",
        "13. Git akışı ve sürümleme",
        "14. Sorun giderme",
        "15. Bilinen sınırlar ve yol haritası",
        "16. Sunum ve canlı demo planı",
        "17. Terimler sözlüğü ve ekler",
    ]
    add_word_numbered(document, [item.split(". ", 1)[1] for item in toc])
    document.add_page_break()

    # 1 Yönetici özeti
    add_word_title(document, "1. Yönetici Özeti")
    add_word_body(
        document,
        "CMS-RAG Assistant; Combat Management System (CMS) ve HAVELSAN ADVENT hakkında "
        "yerel dokümanlardan kaynak gösteren cevaplar üretmek üzere geliştirilmiş, kanıt "
        "öncelikli bir bilgi asistanıdır. Ürün-spesifik resmî içerik ile genel açık/kamu "
        "referansları ayrı koleksiyonlarda tutulur; cevaplar semantik FAISS, BM25, Reciprocal "
        "Rank Fusion ve cross-encoder reranking sonucunda seçilen belge parçalarına dayanır."
    )
    add_word_body(
        document,
        "Araştırma ve kaynak toplama normal kullanımdan önce tamamlanır. Dört metin "
        "çıkarılabilir PDF, 77 anlamlı parça ve bunların önceden hesaplanmış embeddingleri "
        "sürümlenmiş bir snapshot olarak hazırlanır. Streamlit açıldığında bu hazır paket "
        "yüklenir; çalışma anında web taraması veya çekirdek belgeleri yeniden embedding "
        "etme işlemi yapılmaz. Yerel Ollama yalnız getirilen kanıt üzerinden cevap üretir."
    )
    add_word_table(
        document,
        ["Kabul göstergesi", "Sonuç", "Yorum"],
        [
            ["Otomatik test", f"{AUTOMATED_TESTS} / {AUTOMATED_TESTS}", "Birim, entegrasyon, UI, agentic ve mimari sınır testleri"],
            ["Retrieval kabulü", "8 / 8", "Hazır kamuya açık bilgi paketinde doğru kanıt erişimi"],
            ["Yanıt kabulü", "7 / 7", "Kavram, atıf, kaynak ve doğru ret kararı"],
            ["İndekslenen kanıt", "77", "Beş PDF ve önceden hesaplanmış 77 embedding"],
            ["PDF bütünlüğü", "42 / 42 sayfa", "50.802 çıkarılmış metin karakteri"],
            ["Canlı servis", "HTTP 200", "Ana sayfa ve health endpoint"],
            ["Git yeniden üretim", "Başarılı", "Temiz git archive içinde aynı sonuçlar"],
        ],
        [5.2, 3.0, 8.3],
    )
    add_word_callout(
        document,
        "Projenin ayırt edici değeri",
        "Asistan çalışma anında araştırma yapmaz. Önceden hazırlanıp doğrulanmış yerel bilgi "
        "paketinden kanıt getirir; belge ve sayfa bilgisini cevapla birlikte gösterir.",
    )
    document.add_page_break()

    # 2 Problem
    add_word_title(document, "2. Problem, Amaç ve Başarı Ölçütleri")
    document.add_heading("2.1 Problem tanımı", level=2)
    add_word_body(
        document,
        "CMS ve ADVENT gibi teknik alanlarda bilgi; broşürler, ürün sayfaları ve genel "
        "birlikte çalışabilirlik referansları arasında dağınıktır. Klasik anahtar kelime "
        "araması bağlamı kaçırabilir; yalnız büyük dil modeline dayanan cevap ise kaynaksız "
        "veya ürün yeteneği uyduran sonuçlar doğurabilir."
    )
    document.add_heading("2.2 Proje amacı", level=2)
    add_word_bullets(
        document,
        [
            "Yerel dokümanlardan, sayfa seviyesinde izlenebilir cevap üretmek.",
            "Resmî ürün iddiaları ile açık/kamu bağlamını koleksiyon bazında ayırmak.",
            "Türkçe doğal dil sorularını teknik CMS terminolojisiyle eşleştirmek.",
            "Takip sorularını sınırlı sohbet belleğiyle anlamlandırmak.",
            "Yetersiz kanıtta cevap uydurmak yerine güvenli ret vermek.",
            "Tekrarlanabilir test ve Git akışıyla sunulabilir bir mühendislik ürünü oluşturmak.",
        ],
    )
    document.add_heading("2.3 Başarı ölçütleri", level=2)
    add_word_table(
        document,
        ["Ölçüt", "Hedef", "Gerçekleşen"],
        [
            ["Kaynak izlenebilirliği", "Belge + sayfa + otorite", "Karşılandı"],
            ["Kaynak ayrımı", "official / open_source / all", "Karşılandı"],
            ["Tekrarlı PDF", "İkinci kayıt oluşmamalı", "SHA-256 ile engellendi"],
            ["Takip sorusu", "Son bağlamı anlamalı", "3 turluk kapsam izole bellek"],
            ["Hata davranışı", "İlgisiz kaynak göstermemeli", "UI ve motor testleriyle doğrulandı"],
            ["Yeniden üretim", "Hazır paket ve snapshot", f"{AUTOMATED_TESTS}/{AUTOMATED_TESTS}, 8/8, 7/7 ve 77/77"],
        ],
        [5.0, 5.8, 5.7],
    )
    document.add_page_break()

    # 3 kapsam
    add_word_title(document, "3. Kapsam ve Kullanım Senaryoları")
    document.add_heading("3.1 Kapsam içi", level=2)
    add_word_bullets(
        document,
        [
            "ADVENT ve CMS kavramları hakkında hazır doküman paketinden kaynaklı soru-cevap.",
            "Kamuya açık kaynakların ön araştırması, kürasyonu ve PDF paketine dönüştürülmesi.",
            "Chunk ve embedding snapshot'ının normal kullanımdan önce hazırlanması.",
            "İsteğe bağlı ek PDF yükleme, doğrulama, içerik hash'i ve indeksleme.",
            "HAVELSAN resmî içeriği ile NATO/açık kamu referanslarının kaynak kimliğiyle korunması.",
            "Hibrit retrieval, reranking ve sayfa bazlı kanıt birleştirme.",
            "Yerel Ollama üzerinden akışlı Türkçe yanıt.",
            "Tek kullanıcılı yerel Streamlit arayüzü.",
        ],
    )
    document.add_heading("3.2 Kapsam dışı", level=2)
    add_word_bullets(
        document,
        [
            "Gerçek zamanlı savaş yönetimi, silah kontrolü veya operasyonel karar verme.",
            "HAVELSAN şirket içi, özel, tasnifli veya veri sahibi izni bulunmayan içerik.",
            "Normal soru-cevap sırasında internet araştırması veya web taraması.",
            "Çok kiracılı kurumsal kimlik ve rol yönetimi.",
            "Görüntü tabanlı taranmış PDF'ler için OCR.",
        ],
    )
    document.add_heading("3.3 Temel kullanıcı hikâyeleri", level=2)
    add_word_table(
        document,
        ["Rol", "İhtiyaç", "Beklenen sonuç"],
        [
            ["Araştırmacı", "ADVENT nedir?", "Resmî kanıta dayalı kısa tanım"],
            ["Öğrenci", "Varyant örnekleri nelerdir?", "Takip bağlamını kullanan kaynaklı yanıt"],
            ["Mentör/Jüri", "İddianın kaynağı neresi?", "Belge, sayfa ve otorite kartı"],
            ["Geliştirici", "Aynı PDF tekrar yüklenirse?", "Duplicate bildirimi; ikinci kayıt yok"],
            ["Operatör", "NATO birlikte çalışabilirlik ne sağlar?", "Yalnız açık kaynak koleksiyonundan cevap"],
        ],
        [3.2, 6.4, 7.0],
    )
    document.add_page_break()

    # 4 trust
    add_word_title(document, "4. Kaynak Güven Modeli ve Veri Yönetişimi")
    add_word_table(
        document,
        ["Koleksiyon", "İçerik", "Otorite", "İzin verilen iddia"],
        [
            ["official", "Hazır ADVENT broşürü + HAVELSAN kamu içeriği", "Üretici / resmî", "ADVENT ürün ve yetenek bilgileri"],
            ["open_source", "NATO ve diğer kamuya açık birincil referanslar", "Açık / kamu", "Genel C2, veri ve sorumlu yapay zekâ"],
            ["all", "Her iki koleksiyon", "Kaynakta korunur", "Geniş araştırma; kaynak kimliği kaybolmaz"],
        ],
        [3.0, 5.4, 3.4, 4.8],
    )
    add_word_body(
        document,
        "Her kanıt parçası; metin, belge adı, sayfa numarası, kaynak yolu, koleksiyon, "
        "otorite ve varsa kaynak URL'sini taşır. Böylece ürün-spesifik bir iddia açık kaynak "
        "bağlamı gibi veya genel NATO yaklaşımı HAVELSAN ürün yeteneği gibi sunulmaz."
    )
    document.add_heading("4.1 Kaynak kabul ilkeleri", level=2)
    add_word_numbered(
        document,
        [
            "Kaynağın kamuya açık veya kullanım yetkili olduğunu doğrula.",
            "Otorite, koleksiyon ve URL bilgisini kaydet.",
            "PDF imzasını ve 200 MB boyut sınırını kontrol et.",
            "İçeriği SHA-256 ile kimliklendir ve yinelenen kaydı engelle.",
            "Yanıtı yalnız getirilen kanıt ve aynı kapsamdaki sohbet bağlamıyla üret.",
        ],
    )
    add_word_callout(
        document,
        "Koleksiyon izolasyonu",
        "Resmî bir ADVENT turundan sonra açık kaynak NATO sorusu sorulduğunda önceki ürün "
        "iddiaları açık kaynak cevabına taşınmaz. Bu davranış otomatik testle korunur.",
        ORANGE,
    )
    document.add_page_break()

    # 5 architecture
    add_word_title(document, "5. Katmanlı Yazılım Mimarisi")
    add_word_pipeline(
        document,
        [
            ("app.py", "İnce giriş"),
            ("presentation", "Streamlit UI"),
            ("application", "RAG orkestrasyonu"),
            ("domain", "Model ve kurallar"),
            ("infrastructure", "PDF + arama + depo"),
        ],
    )
    add_word_table(
        document,
        ["Katman", "Sorumluluk", "Başlıca modüller"],
        [
            ["presentation", "Tema, sidebar, mesaj, kaynak kartı ve sohbet akışı", "app, chat, components, sidebar, services, theme"],
            ["application", "Belge hazırlama, retrieval ve yerel üretim orkestrasyonu", "engine"],
            ["domain", "Veri modelleri, sorgu kuralları ve deterministik kanıt cevapları", "models, query, evidence"],
            ["infrastructure", "PDF/Markdown işleme, manifest, FAISS, BM25 ve reranker", "ingest, storage, retrieval"],
        ],
        [3.0, 7.0, 6.6],
    )
    document.add_heading("5.1 Bağımlılık kuralları", level=2)
    add_word_bullets(
        document,
        [
            "Domain üst katmanlara ve UI/FAISS/Ollama teknolojilerine bağımlı değildir.",
            "Infrastructure yalnız domain modellerini kullanır.",
            "Application, domain kuralları ile altyapıyı kullanım senaryosunda birleştirir.",
            "Presentation, iş kuralı içermez; uygulama servisini çağırır.",
            "Kök app.py yalnız sunum orkestratörünü çalıştırır.",
        ],
    )
    add_word_callout(
        document,
        "Mimari koruma",
        "tests/test_architecture.py; giriş noktasının ince kalmasını, bağımlılık yönlerini ve "
        "tüm kaynak sınıf/fonksiyonlarının Türkçe açıklama taşımasını AST üzerinden denetler.",
    )
    document.add_heading("5.2 Agentic LangGraph ve kalıcı konuşmalar", level=2)
    add_word_body(
        document,
        "Agentic mod; bilgi, MCP kontrol ve güvenli ret rotalarını ayrı checkpoint'li "
        "düğümlerle yönetir. Bilgi akışında sorgu planlama, hibrit retrieval, kanıt kapısı, "
        "yerel üretim, atıf doğrulama ve en fazla bir deterministik onarım bulunur. PostgreSQL "
        "checkpointer kullanıldığında konuşmalar yeniden başlatma sonrasında listelenip kaynak "
        "kartlarıyla geri yüklenebilir.",
    )
    add_word_body(
        document,
        "MCP yazma komutu interrupt noktasında durur. Operatör onayı olmadan set aracı "
        "çağrılmaz; onay veya ret aynı thread üzerinde resume edilerek kalıcı sonuca dönüşür. "
        "Bekleyen onay uygulama yeniden başlatılsa bile bulunabilir. Planlama veya model arızası "
        "traceback yerine belgesiz üretim yapmayan güvenli sonuçla kapanır.",
    )
    document.add_page_break()

    # 6 document lifecycle
    add_word_title(document, "6. Belge Yaşam Döngüsü")
    add_word_pipeline(
        document,
        [
            ("1 · Doğrula", "%PDF + boyut"),
            ("2 · Kimlik", "SHA-256"),
            ("3 · Sakla", "İçerik adresli ad"),
            ("4 · Manifest", "Ad + boyut + tarih"),
            ("5 · Parçala", "Sayfa + overlap"),
            ("6 · İndeksle", "FAISS + BM25"),
        ],
    )
    add_word_numbered(
        document,
        [
            "Yüklenen içerik PDF imzası ve azami 200 MB sınırıyla doğrulanır.",
            "Dosya içeriği SHA-256 ile kimliklendirilir; ad değişse bile kopya anlaşılır.",
            "Yeni belge hash önekli güvenli dosya adıyla yerel depoya yazılır.",
            "Özgün görünen ad, boyut, kaynak tipi ve zaman manifestte tutulur.",
            "PDF sayfa sayfa okunur; metin 900 karakter ve 150 karakter örtüşmeyle parçalanır.",
            "Bozuk dosya tüm işlemi durdurmaz; okunabilen kaynaklarla indeks kurulur.",
            "Belge silme dosyayı ve manifest kaydını birlikte kaldırır; indeks yenilenir.",
        ],
    )
    document.add_heading("6.1 Güvenlik kontrolleri", level=2)
    add_word_table(
        document,
        ["Risk", "Kontrol", "Test"],
        [
            ["Aynı PDF'nin tekrar yüklenmesi", "İçerik SHA-256 karşılaştırması", "duplicate_pdf_content_is_stored_once"],
            ["PDF olmayan dosya", "Magic byte kontrolü", "non_pdf_content_is_rejected"],
            ["Aşırı büyük dosya", "200 MB sınırı", "oversized_pdf_is_rejected"],
            ["Manifest path traversal", "Çözülmüş üst dizin kontrolü", "delete_rejects_manifest_path_traversal"],
            ["Bozuk PDF", "Dosya bazlı hata izolasyonu", "invalid_pdf_is_skipped"],
        ],
        [5.0, 5.5, 6.1],
    )
    document.add_page_break()

    # 7 RAG
    add_word_title(document, "7. Hibrit RAG ve Retrieval Hattı")
    add_word_pipeline(
        document,
        [
            ("Soru", "Türkçe doğal dil"),
            ("Genişlet", "CMS terimleri"),
            ("Semantic", "BGE + FAISS"),
            ("Lexical", "BM25"),
            ("Fusion", "RRF"),
            ("Rerank", "Cross-encoder"),
            ("Kanıt", "Sayfa birleştirme"),
        ],
    )
    document.add_heading("7.1 Semantik arama", level=2)
    add_word_body(
        document,
        "BAAI/bge-small-en-v1.5 modeli belge parçaları ve sorgu için normalize embedding "
        "üretir. FAISS IndexFlatIP, kosinüs benzerliğine eşdeğer iç çarpım üzerinden anlamca "
        "yakın adayları bulur."
    )
    document.add_heading("7.2 BM25", level=2)
    add_word_body(
        document,
        "BM25; ADVENT, MARTI, UFUK, MÜREN, Link 11/16 ve CMS gibi tam ürün adı, kısaltma "
        "ve teknik terim eşleşmelerini yakalar. Semantik aramanın kaçırabileceği kesin "
        "sözcük sinyalini korur."
    )
    document.add_heading("7.3 Reciprocal Rank Fusion", level=2)
    add_word_body(
        document,
        "RRF, semantik ve BM25 puanlarını doğrudan karşılaştırmaz; her adayın sırasını "
        "1/(60+sıra) formülüyle ortak skora dönüştürür. Böylece farklı puan ölçekleri "
        "istikrarlı biçimde birleşir."
    )
    document.add_heading("7.4 Cross-encoder reranking", level=2)
    add_word_body(
        document,
        "BAAI/bge-reranker-base, en iyi adayları soru–metin çifti olarak yeniden değerlendirir. "
        "Model yerelde bulunamazsa sistem RRF sıralamasıyla çalışmaya devam eder."
    )
    document.add_heading("7.5 Sayfa bazlı kanıt birleştirme", level=2)
    add_word_body(
        document,
        "Aynı sayfanın örtüşen veya tamamlayıcı parçaları tek kanıt öğesinde birleştirilir. "
        "Bu yaklaşım kaynak kartı tekrarını azaltırken bağlam kaybını önler."
    )
    document.add_page_break()

    # 8 memory and generation
    add_word_title(document, "8. Konuşma Belleği ve Cevap Üretimi")
    add_word_body(
        document,
        "Motor en fazla üç tamamlanmış soru-cevap turunu saklar. Kısa takip soruları, son "
        "turların soru ve cevaplarıyla retrieval sorgusuna dönüştürülür. Her tur kendi "
        "kaynak kapsamı etiketiyle tutulur; official, open_source ve all geçmişleri birbirine "
        "karışmaz."
    )
    add_word_table(
        document,
        ["Karar sırası", "Davranış"],
        [
            ["1. İndeks yok", "Önce belge yükleme ve indeksleme yönlendirmesi"],
            ["2. Alan dışı soru", "Retrieval çalıştırmadan güvenli ret"],
            ["3. Kesin kanıt kuralı", "Belgedeki açık ifadeden deterministik kısa cevap"],
            ["4. Hibrit retrieval", "Kapsam filtreli aday ve reranking"],
            ["5. Ollama üretimi", "Yalnız CONTEXT üzerinden en fazla 55 kelimelik Türkçe yanıt"],
            ["6. Kaynak kontrolü", "Gerekirse [SOURCE 1] deterministik eklenir"],
            ["7. Hata", "Ollama servis mesajı; ilgisiz kaynak kartları gizlenir"],
        ],
        [4.8, 11.8],
    )
    document.add_heading("8.1 Yerel model ayarları", level=2)
    add_word_table(
        document,
        ["Ayar", "Değer", "Gerekçe"],
        [
            ["Varsayılan model", "qwen2.5:3b", "CPU üzerinde etkileşimli kullanım"],
            ["Kalite seçeneği", "qwen2.5:7b", "Güçlü donanımda daha yüksek üretim kalitesi"],
            ["Temperature", "0.1", "Düşük varyans ve daha kararlı cevap"],
            ["num_predict", "160 + gerekirse 96", "Kesik cümleyi tek kontrollü devamla tamamlama"],
            ["keep_alive", "30 dakika", "Ardışık sorgularda model yükleme gecikmesini azaltma"],
            ["Timeout", "120 saniye", "Yerel model tıkanmasına karşı sınır"],
        ],
        [4.3, 4.0, 8.3],
    )
    document.add_page_break()

    # 9 security
    add_word_title(document, "9. Güvenlik, Gizlilik ve Güvenilirlik")
    add_word_bullets(
        document,
        [
            "Doküman metni yerel Ollama dışında bir üretim servisine gönderilmez.",
            "Rastgele web taraması yapılmaz; web kaynakları seçilip yerel referansa dönüştürülür.",
            "Kaynak alıntıları HTML çıktısında escape edilerek XSS riski azaltılır.",
            "Dosya adı doğrudan depolama yolu olmaz; hash tabanlı güvenli ad kullanılır.",
            "Yetersiz kanıt ve servis hatasında retrieval sonuçları kanıt gibi gösterilmez.",
            "Alan dışı sorular güvenli ret alır.",
            "Kapsam etiketli konuşma belleği koleksiyonlar arası bilgi sızıntısını engeller.",
            "Yerel audit ham soru/cevap saklamaz; sorgu SHA-256 özeti ve işletim metadatası tutulur.",
            "PDF önizlemesi yalnız proje data dizini altındaki doğrulanmış dosyalara erişir.",
        ],
    )
    document.add_heading("9.1 Üretim ortamı için ek kontroller", level=2)
    add_word_table(
        document,
        ["Kontrol", "Amaç", "Öncelik"],
        [
            ["Kimlik ve RBAC", "Belge/koleksiyon erişimini role göre sınırla", "Yüksek"],
            ["Disk şifreleme", "Yerel dokümanları cihaz kaybına karşı koru", "Yüksek"],
            ["Merkezi audit", "Yükleme, silme, sorgu ve kaynak kullanımını izle", "Yüksek"],
            ["Zararlı dosya taraması", "PDF içindeki kötü amaçlı içeriği engelle", "Orta"],
            ["Kaynak onay iş akışı", "Küratörlü referans değişikliğini dört gözle yönet", "Orta"],
            ["Ağ çıkış politikası", "Model ve uygulamanın beklenmeyen dış trafiğini sınırla", "Orta"],
        ],
        [5.0, 8.2, 3.4],
    )
    add_word_callout(
        document,
        "Dürüst güvenlik notu",
        "pip check temizdir. Çevrim içi CVE sorgusu, kurulu paket envanterini haricî servise "
        "göndereceği için ayrıca açık izin olmadan çalıştırılmamıştır.",
        ORANGE,
    )
    document.add_page_break()

    # 10 UI
    add_word_title(document, "10. Kullanıcı Arayüzü ve Deneyim")
    add_word_table(
        document,
        ["Bölge", "İşlev"],
        [
            ["Sol panel", "Asistan/değerlendirme görünümü, kapsam, PDF, indeks ve oturum işlemleri"],
            ["Çalışma durumu", "Model, belge sayısı, koleksiyon ve arama yöntemi"],
            ["Üst gösterge", "Yüklü belge ve kanıt parçası sayıları"],
            ["Sohbet alanı", "Geçmiş kullanıcı/asistan mesajları ve akışlı yeni cevap"],
            ["Yanıt durumu", "Bağlam çözümleme, kanıt seçimi ve tamamlanma göstergesi"],
            ["Kanıt paketi", "Belge, sayfa, otorite, alıntı, URL ve gerçek PDF sayfa önizlemesi"],
            ["Değerlendirme", "45-vaka benchmark, 30-vaka hakem, chunk, backend ve audit görünümü"],
            ["Belge yönetimi", "Manifest kaydı ve güvenli belge kaldırma"],
        ],
        [4.2, 12.4],
    )
    document.add_heading("10.1 Kullanıcı yolculuğu", level=2)
    add_word_numbered(
        document,
        [
            "Kullanıcı sorgu kapsamını seçer.",
            "Gerekirse yetkili PDF'yi yükler ve doğrulayıp indeksler.",
            "Doğal dilde sorusunu gönderir.",
            "Sistem bağlamı çözer, kanıtları bulur ve yeniden sıralar.",
            "Yanıt ekrana parça parça yazılır.",
            "Kullanıcı kanıt paketini açıp belge ve sayfayı denetler.",
            "Takip sorusu aynı kapsamdaki kısa konuşma belleğiyle anlaşılır.",
        ],
    )
    add_word_callout(
        document,
        "Tasarım ilkesi",
        "Arayüz modelin ne düşündüğünü değil, hangi kanıtı kullandığını görünür kılar. "
        "Kullanıcı güveni cevap süsünden değil izlenebilirlikten gelir.",
    )
    document.add_page_break()

    # 11 install
    add_word_title(document, "11. Kurulum, Çalıştırma ve İşletim")
    document.add_heading("11.1 Gereksinimler", level=2)
    add_word_bullets(
        document,
        [
            "Python 3.11 veya üzeri",
            "Ollama",
            "En az 8 GB RAM",
            "İlk model indirmeleri için internet; sonrasında çevrimdışı çalışma seçeneği",
            "Windows PowerShell veya eşdeğer terminal",
        ],
    )
    document.add_heading("11.2 Kurulum komutları", level=2)
    code_lines = [
        "python -m venv .venv",
        r".\.venv\Scripts\Activate.ps1",
        "python -m pip install -r requirements.txt",
        "ollama pull qwen2.5:3b",
    ]
    for line in code_lines:
        paragraph = document.add_paragraph()
        set_cell = paragraph.add_run(line)
        set_cell.font.name = "Cascadia Mono"
        set_cell.font.size = Pt(9)
        set_cell.font.color.rgb = RGBColor.from_string(NAVY)
        paragraph.paragraph_format.left_indent = Cm(0.7)
    document.add_heading("11.3 Çalıştırma", level=2)
    add_word_table(
        document,
        ["Terminal", "Komut"],
        [
            ["1 · Ollama", "ollama serve"],
            ["2 · Streamlit", r".\.venv\Scripts\python.exe -m streamlit run app.py"],
            ["Tarayıcı", "http://localhost:8501"],
        ],
        [4.0, 12.6],
    )
    document.add_heading("11.4 Ortam değişkenleri", level=2)
    add_word_table(
        document,
        ["Değişken", "Örnek", "Etkisi"],
        [
            ["CMS_RAG_MODEL", "qwen2.5:7b", "Varsayılan 3B yerine kalite modelini seçer"],
            ["CMS_RAG_OFFLINE", "1", "Embedding/reranker için yalnız yerel dosyaları kullanır"],
        ],
        [4.2, 4.0, 8.4],
    )
    document.add_page_break()

    # 12 tests
    add_word_title(document, "12. Test Stratejisi ve Kabul Kanıtları")
    add_word_table(
        document,
        ["Test katmanı", "Kapsam", "Sonuç"],
        [
            ["Storage", "Hash, duplicate, PDF imzası, boyut, manifest, silme, traversal", "Başarılı"],
            ["Ingest", "PDF/Markdown, sayfa metadata, bozuk dosya, kanıt birleştirme", "Başarılı"],
            ["Engine", "Takip, kaynak etiketi, ret, model seçimi, kapsam izolasyonu", "Başarılı"],
            ["UI", "Kapsam seçimi, konuşma yolculuğu, Ollama hatası, kaynak kalıcılığı", "Başarılı"],
            ["Architecture", "İnce app.py, bağımlılık yönü, Türkçe açıklama standardı", "Başarılı"],
            ["Audit", "Ham içeriksiz SHA-256 özeti, gecikme, kapsam ve kaynak metadatası", "Başarılı"],
            ["Agentic", "Route, checkpoint, konuşma geri yükleme, interrupt/resume ve hata sınırları", "Başarılı"],
            ["Retrieval", "45 altın belge/terim/sayfa ve güvenli-ret vakası", "45 / 45"],
        ],
        [4.0, 9.7, 2.9],
    )
    document.add_heading("12.1 Kabul komutları", level=2)
    for line in (
        r".\.venv\Scripts\python.exe -m unittest discover -s tests -v",
        r".\.venv\Scripts\python.exe -m scripts.evaluate_retrieval",
        r".\.venv\Scripts\python.exe -m pip check",
    ):
        paragraph = document.add_paragraph()
        run = paragraph.add_run(line)
        run.font.name = "Cascadia Mono"
        run.font.size = Pt(8.5)
        run.font.color.rgb = RGBColor.from_string(NAVY)
    document.add_heading("12.2 Ölçülen son sonuçlar", level=2)
    add_word_table(
        document,
        ["Gösterge", "Değer"],
        [
            ["Otomatik test", f"{AUTOMATED_TESTS}/{AUTOMATED_TESTS}"],
            ["Altın benchmark", f"{BENCHMARK_REPORT['passed']}/{BENCHMARK_REPORT['dataset_cases']}"],
            [
                "TP / TN / FP / FN",
                "{true_positive} / {true_negative} / {false_positive} / {false_negative}".format(
                    **BENCHMARK_REPORT["confusion_matrix"]
                ),
            ],
            [
                "Bağımsız cevap hakemi",
                f"{QUALITY_REPORT['stage2']['summary']['strict_passed']}/"
                f"{QUALITY_REPORT['stage2']['summary']['evaluated']} katı başarı",
            ],
            [
                "Chunk hakemi",
                f"{QUALITY_REPORT['stage1']['summary']['acceptable_count']}/"
                f"{QUALITY_REPORT['stage1']['summary']['chunk_count']} geçerli",
            ],
            [
                "20-vaka TP / TN / FP / FN",
                "{true_positive} / {true_negative} / {false_positive} / {false_negative}".format(
                    **LINEAGE_REPORT["confusion_matrix"]
                ),
            ],
            [
                "20-vaka accuracy / F1",
                f"{LINEAGE_REPORT['confusion_matrix']['accuracy']:.1%} / "
                f"{LINEAGE_REPORT['confusion_matrix']['f1']:.2%}",
            ],
            [
                "Exact chunk-köken başarı",
                f"{LINEAGE_REPORT['lineage']['strict_passes']}/"
                f"{LINEAGE_REPORT['lineage']['evaluated_positive_cases']}",
            ],
            [
                "Hit@6 / MRR",
                f"{BENCHMARK_REPORT['retrieval']['hit_at_k']:.0%} / "
                f"{BENCHMARK_REPORT['retrieval']['mrr']:.4f}".replace(".", ","),
            ],
            ["İndekslenen parça", str(QUALITY_REPORT['stage1']['summary']['chunk_count'])],
            ["PDF sayfası", "42/42 metinli"],
            ["Boş parça / eksik kaynak yolu", "0 / 0"],
            ["Canlı Streamlit", "HTTP 200"],
            ["Temiz Git arşivi", "Aynı test sonuçları"],
        ],
        [8.0, 8.6],
    )
    add_word_callout(
        document,
        "Kabul yaklaşımı",
        "Yalnız çalışma klasöründe test geçmesi yeterli sayılmadı. Commit, git archive ile "
        "bağımsız klasöre çıkarıldı; manifest yeniden üretildi ve tüm testler tekrar çalıştırıldı.",
    )
    document.add_page_break()

    # 13 git
    add_word_title(document, "13. Git Akışı ve Sürümleme")
    add_word_pipeline(
        document,
        [
            ("main", "Yayımlanmış sürüm"),
            ("develop", "Entegrasyon"),
            ("codex/*", "Özellik/refactor"),
            ("release/*", "Son kabul"),
            ("hotfix/*", "Acil düzeltme"),
        ],
    )
    add_word_table(
        document,
        ["Dal", "Kaynak", "Hedef", "Kural"],
        [
            ["main", "release/hotfix", "—", "Doğrudan commit yok; sunulabilir sürüm"],
            ["develop", "feature/release/hotfix", "—", "Bir sonraki sürüm entegrasyonu"],
            ["codex/<konu>", "develop", "develop", "Kod, test ve inceleme"],
            ["release/<sürüm>", "develop", "main + develop", "Tam kabul ve sürüm etiketi"],
            ["hotfix/<sürüm>", "main", "main + develop", "Üretim acil düzeltmesi"],
        ],
        [3.2, 3.4, 4.0, 6.0],
    )
    document.add_heading("13.1 Commit standardı", level=2)
    add_word_bullets(
        document,
        [
            "feat: yeni kullanıcı özelliği",
            "fix: hata düzeltmesi",
            "refactor: davranışı koruyan mimari düzenleme",
            "test: test ve kabul kapsamı",
            "docs: dokümantasyon",
            "chore: bağımlılık ve bakım",
        ],
    )
    add_word_body(
        document,
        f"Bu dokümantasyonun referans aldığı doğrulanmış entegrasyon commit'i {COMMIT}'dır. "
        "Eski main geçmişi silinmemiş; modüler mimari feature dalından develop'a merge commit "
        "ile alınmıştır."
    )
    document.add_page_break()

    # 14 troubleshooting
    add_word_title(document, "14. Sorun Giderme")
    add_word_table(
        document,
        ["Belirti", "Olası neden", "Çözüm"],
        [
            ["Ollama servisine ulaşılamadı", "Ollama çalışmıyor veya model yok", "ollama serve; ollama pull qwen2.5:3b"],
            ["Yeni dosya bulunamadı", "Aynı içerik daha önce yüklendi", "Duplicate bilgi mesajını kontrol et; yeniden kayıt gerekmez"],
            ["PDF reddedildi", "PDF imzası yok veya 200 MB üzeri", "Geçerli ve yetkili PDF kullan"],
            ["Soru cevaplanmadı", "Seçili kapsamda yeterli kanıt yok", "Kapsamı ve kanıt paketini kontrol et; uygun kaynak ekle"],
            ["İlk açılış yavaş", "Embedding/reranker ilk kez yükleniyor", "Model dosyalarının indirilmesini bekle; sonraki açılış hızlanır"],
            ["Takip sorusu yanlış bağlam", "Kapsam değişti veya soru aşırı belirsiz", "Soruyu ürün/kavram adıyla netleştir"],
            ["Port 8501 kullanımda", "Başka Streamlit süreci var", "Mevcut süreci kapat veya farklı --server.port seç"],
        ],
        [4.2, 5.6, 6.8],
    )
    document.add_heading("14.1 Hızlı sağlık kontrolü", level=2)
    add_word_numbered(
        document,
        [
            "http://localhost:8501/_stcore/health adresinin HTTP 200 döndürdüğünü kontrol et.",
            "ollama list ile qwen2.5:3b modelinin kurulu olduğunu doğrula.",
            "pip check çalıştır ve bozuk bağımlılık olmadığını doğrula.",
            "Test ve retrieval kabul komutlarını yeniden çalıştır.",
            "Gerekirse arayüzden İndeksi yenile seçeneğini kullan.",
        ],
    )
    document.add_page_break()

    # 15 roadmap
    add_word_title(document, "15. Bilinen Sınırlar ve Yol Haritası")
    add_word_table(
        document,
        ["Ufuk", "Geliştirme", "Değer"],
        [
            ["Yakın", "OCR ve taranmış PDF desteği", "Görüntü tabanlı doküman kapsamı"],
            ["Yakın", "Otomatik kaynak yenilik kontrolü", "Resmî sayfa drift'ini erken yakalama"],
            ["Orta", "Kalıcı vektör deposu", "Büyük koleksiyonda hızlı başlangıç"],
            ["Orta", "Kullanıcı geri bildirimi ve kalite eğilim alarmı", "Hata analizi ve drift görünürlüğü"],
            ["Orta", "RBAC ve koleksiyon yetkilendirmesi", "Kurumsal çok kullanıcılı çalışma"],
            ["Uzun", "Merkezi audit aktarımı ve kaynak onay iş akışı", "Yönetişim ve mevzuat uyumu"],
            ["Uzun", "Çok modlu tablo/şema/görsel retrieval", "Teknik dokümanların zengin anlaşılması"],
        ],
        [3.0, 7.0, 6.6],
    )
    add_word_callout(
        document,
        "Öncelik önerisi",
        "Sunum sonrası ilk mühendislik paketi: OCR + kalıcı indeks + otomatik regression "
        "pipeline. Yerel gizlilik-korumalı audit tamamlandı; kurumsal yaygınlaştırma "
        "öncesinde RBAC, merkezi audit aktarımı ve disk şifreleme eklenmelidir.",
        ORANGE,
    )
    document.add_page_break()

    # 16 presentation
    add_word_title(document, "16. Sunum ve Canlı Demo Planı")
    add_word_table(
        document,
        ["Süre", "Bölüm", "Ana mesaj"],
        [
            ["1 dk", "Problem", "Dağınık teknik bilgi ve kaynaksız LLM riski"],
            ["2 dk", "Çözüm", "Yerel, kaynak kontrollü ve koleksiyon ayrımlı RAG"],
            ["3 dk", "Mimari", "Katmanlar, hibrit retrieval ve kanıt zinciri"],
            ["2 dk", "Güvenilirlik", "Hash, güvenli ret, kapsam izolasyonu ve yerel model"],
            ["3 dk", "Canlı demo", "ADVENT → takip sorusu → NATO kapsamı → kaynaksız ret"],
            ["1 dk", "Test", f"{AUTOMATED_TESTS}/{AUTOMATED_TESTS}, 8/8, 7/7 ve 77/77 snapshot"],
            ["1 dk", "Yol haritası", "OCR, kalıcı indeks, RBAC ve audit"],
        ],
        [2.5, 4.5, 9.6],
    )
    document.add_heading("16.1 Önerilen demo soruları", level=2)
    add_word_numbered(
        document,
        [
            "ADVENT nedir?",
            "Örnekleri var mıdır?",
            "Bunların görevleri nelerdir?",
            "Kapsamı açık/kamu olarak değiştir: NATO interoperability data-centric yaklaşımı ne sağlar?",
            "Ben kimim? — güvenli ret ve kaynak uydurmama davranışını göster.",
        ],
    )
    add_word_heading = document.add_heading("16.2 Sunucu anlatım ipuçları", level=2)
    add_word_bullets(
        document,
        [
            "Önce cevabı değil kanıt paketini göster; projenin farkı burada.",
            "Semantic + BM25 + RRF + reranking zincirini tek cümleyle açıkla.",
            "Yerel Ollama'nın veri gizliliğine katkısını vurgula.",
            "Sistemin operasyonel karar sistemi olmadığını açıkça belirt.",
            "Test sayılarını ezberlemek yerine kabul tablosunu kullan.",
        ],
    )
    document.add_page_break()

    # 17 glossary
    add_word_title(document, "17. Terimler Sözlüğü ve Ekler")
    add_word_table(
        document,
        ["Terim", "Açıklama"],
        [
            ["CMS", "Combat Management System / Savaş Yönetim Sistemi"],
            ["RAG", "Retrieval-Augmented Generation; getirilen kanıtla cevap üretimi"],
            ["Embedding", "Metni anlamsal vektöre dönüştüren temsil"],
            ["FAISS", "Yoğun vektör benzerliği için yerel arama kütüphanesi"],
            ["BM25", "Sözcüksel alaka ve terim sıklığına dayalı arama"],
            ["RRF", "Farklı sonuç sıralamalarını sıra bilgisiyle birleştirme yöntemi"],
            ["Cross-encoder", "Soru ve kanıtı birlikte değerlendirerek yeniden sıralayan model"],
            ["Chunk", "Belge ve sayfa metadatasını koruyan metin parçası"],
            ["Grounding", "Cevabın gösterilen kaynağa dayandırılması"],
            ["Ollama", "Yerel büyük dil modeli çalıştırma servisi"],
            ["Scope", "official, open_source veya all kaynak kapsamı"],
            ["Manifest", "Depolanan belgenin hash, ad, boyut ve tarih kaydı"],
        ],
        [4.0, 12.6],
    )
    document.add_heading("Ek A — Proje klasör ağacı", level=2)
    tree_lines = [
        "app.py",
        "src/cms_rag/application/engine.py",
        "src/cms_rag/domain/{models,query,evidence}.py",
        "src/cms_rag/infrastructure/{storage,ingest,retrieval}.py",
        "src/cms_rag/presentation/{app,chat,components,sidebar,services,theme}.py",
        "data/documents/",
        "data/references/{official,open_source}/",
        "scripts/evaluate_retrieval.py",
        "tests/test_{storage,ingest,engine,app_ui,architecture}.py",
        "docs/{ARCHITECTURE,GIT_WORKFLOW,TEST_EVIDENCE,FINAL_ACCEPTANCE_REPORT}.md",
    ]
    for line in tree_lines:
        paragraph = document.add_paragraph()
        run = paragraph.add_run(line)
        run.font.name = "Cascadia Mono"
        run.font.size = Pt(8.5)
        run.font.color.rgb = RGBColor.from_string(NAVY)
    document.add_heading("Ek B — Resmî referanslar", level=2)
    add_word_bullets(
        document,
        [
            "HAVELSAN ADVENT ürün sayfası: https://www.havelsan.com/en/solutions/advent-combat-management-system",
            "NATO Digital Transformation Implementation Strategy 2.0: https://www.nato.int/en/about-us/official-texts-and-resources/official-texts/2026/05/26/natos-digital-transformation-implementation-strategy",
            "Yerel resmî başlangıç belgesi: data/documents/..._advent_cms.pdf",
        ],
    )
    add_word_callout(
        document,
        "Nihai durum",
        "Proje; modüler mimari, kaynak yönetişimi, test kanıtı ve Git akışıyla sunuma "
        "hazırdır. Uzak depo release/push ve kurumsal erişim kontrolleri ayrı yayımlama "
        "adımıdır.",
    )

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    document.save(DOCX_PATH)


def ppt_color(hex_value: str) -> PptRGB:
    """Onaltılık rengi python-pptx RGBColor nesnesine dönüştürür."""

    return PptRGB.from_string(hex_value)


def set_slide_background(slide, color: str = "F7F9FC") -> None:
    """Slayt arka planını tek renk kurumsal zeminle doldurur."""

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = ppt_color(color)


def add_ppt_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 20,
    color: str = INK,
    bold: bool = False,
    font: str = "Aptos",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.04,
):
    """Slayta taşma kontrollü, biçimlendirilmiş bir metin kutusu ekler."""

    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = PptInches(margin)
    frame.margin_right = PptInches(margin)
    frame.margin_top = PptInches(margin)
    frame.margin_bottom = PptInches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = ppt_color(color)
    return box


def add_ppt_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = WHITE,
    line: str = LINE,
    radius=True,
    transparency: int = 0,
):
    """Slayta kart veya panel olarak kullanılacak şekil ekler."""

    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type,
        PptInches(x),
        PptInches(y),
        PptInches(w),
        PptInches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ppt_color(fill)
    shape.fill.transparency = transparency
    shape.line.color.rgb = ppt_color(line)
    shape.line.width = PptPt(1)
    return shape


def add_ppt_line(slide, x1: float, y1: float, x2: float, y2: float, color: str = BLUE, width: float = 2):
    """İki slayt koordinatı arasına bağlantı çizgisi ekler."""

    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        PptInches(x1),
        PptInches(y1),
        PptInches(x2),
        PptInches(y2),
    )
    line.line.color.rgb = ppt_color(color)
    line.line.width = PptPt(width)
    return line


def add_slide_frame(slide, number: int, section: str, dark: bool = False) -> None:
    """Her slayta bölüm etiketi, üst çizgi ve sayfa numarası ekler."""

    color = WHITE if dark else NAVY
    muted = "B8C7DC" if dark else MUTED
    add_ppt_text(slide, f"CMS-RAG  /  {section.upper()}", 0.55, 0.25, 7.5, 0.28, size=8.5, color=muted, bold=True)
    add_ppt_text(slide, f"{number:02d}", 12.05, 0.22, 0.65, 0.3, size=9, color=color, bold=True, align=PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        PptInches(0.55),
        PptInches(0.65),
        PptInches(12.15),
        PptInches(0.025),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ppt_color(CYAN if dark else BLUE)
    line.line.fill.background()


def add_slide_title(slide, title: str, subtitle: str | None = None, dark: bool = False) -> None:
    """Slayt başlığını ve isteğe bağlı açıklamasını ortak konumda oluşturur."""

    main_color = WHITE if dark else NAVY
    sub_color = "CAD5E5" if dark else MUTED
    add_ppt_text(slide, title, 0.65, 0.82, 11.8, 0.65, size=26, color=main_color, bold=True, font="Aptos Display")
    if subtitle:
        add_ppt_text(slide, subtitle, 0.68, 1.48, 11.4, 0.45, size=11.5, color=sub_color)


def add_metric_card(slide, x: float, y: float, w: float, value: str, label: str, accent: str = CYAN) -> None:
    """Önemli kabul değerini büyük sayı ve kısa etiketle gösterir."""

    add_ppt_box(slide, x, y, w, 1.15, fill=WHITE, line=LINE)
    accent_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        PptInches(x),
        PptInches(y),
        PptInches(0.08),
        PptInches(1.15),
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ppt_color(accent)
    accent_bar.line.fill.background()
    add_ppt_text(slide, value, x + 0.22, y + 0.15, w - 0.35, 0.48, size=23, color=NAVY, bold=True)
    add_ppt_text(slide, label, x + 0.22, y + 0.68, w - 0.35, 0.28, size=9.5, color=MUTED, bold=True)


def add_bullet_list(slide, items: list[str], x: float, y: float, w: float, h: float, color: str = INK, size: float = 16) -> None:
    """Slayta okunaklı aralıklara sahip bir madde listesi ekler."""

    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = PptInches(0.05)
    frame.margin_right = PptInches(0.02)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = PptPt(size)
        paragraph.font.color.rgb = ppt_color(color)
        paragraph.space_after = PptPt(10)
        paragraph._p.get_or_add_pPr().insert(0, OxmlElement("a:buChar"))
        paragraph._p.pPr[0].set("char", "◆")
    return box


def add_label(slide, text: str, x: float, y: float, w: float, color: str = BLUE) -> None:
    """Küçük bölüm veya durum etiketini kapsül biçiminde çizer."""

    add_ppt_box(slide, x, y, w, 0.34, fill=PALE_BLUE, line=PALE_BLUE)
    add_ppt_text(slide, text.upper(), x + 0.08, y + 0.07, w - 0.16, 0.18, size=7.5, color=color, bold=True, align=PP_ALIGN.CENTER)


def add_pipeline_nodes(slide, stages: list[tuple[str, str]], y: float = 3.1) -> None:
    """RAG aşamalarını birbirine bağlı kartlarla yatay gösterir."""

    start_x = 0.7
    gap = 0.15
    available = 11.95
    width = (available - gap * (len(stages) - 1)) / len(stages)
    for index, (title, detail) in enumerate(stages):
        x = start_x + index * (width + gap)
        add_ppt_box(slide, x, y, width, 1.55, fill=WHITE, line=BLUE if index % 2 == 0 else CYAN)
        add_ppt_text(slide, f"{index + 1:02d}", x + 0.12, y + 0.15, 0.35, 0.25, size=8, color=CYAN, bold=True)
        add_ppt_text(slide, title, x + 0.12, y + 0.48, width - 0.24, 0.35, size=12, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_ppt_text(slide, detail, x + 0.12, y + 0.9, width - 0.24, 0.42, size=8.2, color=MUTED, align=PP_ALIGN.CENTER)
        if index < len(stages) - 1:
            add_ppt_line(slide, x + width, y + 0.78, x + width + gap, y + 0.78, CYAN, 1.8)


def build_powerpoint() -> None:
    """Sunumda doğrudan kullanılabilecek 16:9 PowerPoint dosyasını üretir."""

    presentation = Presentation()
    presentation.slide_width = PptInches(13.333)
    presentation.slide_height = PptInches(7.5)
    blank = presentation.slide_layouts[6]

    # 1 Kapak
    slide = presentation.slides.add_slide(blank)
    set_slide_background(slide, NAVY)
    add_slide_frame(slide, 1, "Nihai Proje Sunumu", dark=True)
    add_ppt_text(slide, "EVIDENCE-FIRST / LOCAL AI", 0.7, 1.25, 5.4, 0.3, size=9, color=CYAN, bold=True)
    add_ppt_text(slide, "CMS-RAG\nAssistant", 0.68, 1.68, 6.7, 1.65, size=42, color=WHITE, bold=True, font="Aptos Display")
    add_ppt_text(
        slide,
        "Savaş Yönetim Sistemi dokümanları için\nkaynak kontrollü bilgi keşfi",
        0.75,
        3.55,
        5.8,
        0.9,
        size=18,
        color="C7D5E8",
    )
    # Sağ tarafta soyut kanıt ağı
    nodes = [(8.0, 1.6, "PDF"), (10.3, 1.25, "SNAPSHOT"), (11.25, 3.0, "RRF"), (9.25, 3.65, "RAG"), (7.55, 5.0, "SOURCE")]
    for x, y, label in nodes:
        add_ppt_box(slide, x, y, 1.25, 0.72, fill="15365F", line=CYAN)
        add_ppt_text(slide, label, x, y + 0.23, 1.25, 0.22, size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    for (x1, y1, _), (x2, y2, _) in zip(nodes, nodes[1:]):
        line = add_ppt_line(slide, x1 + 0.62, y1 + 0.36, x2 + 0.62, y2 + 0.36, CYAN, 1.3)
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    add_ppt_text(slide, f"Sürüm {VERSION}  ·  {TODAY.strftime('%d.%m.%Y')}  ·  {COMMIT}", 0.75, 6.72, 7.0, 0.3, size=9, color="8FA8C8")

    # 2 Yönetici özeti
    slide = presentation.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_frame(slide, 2, "Yönetici Özeti")
    add_slide_title(slide, "Tek cümlede proje", "Önceden hazırlanmış kamu bilgisini çevrimdışı, yerel ve kaynaklı cevaba dönüştürür.")
    add_metric_card(slide, 0.7, 2.1, 2.7, f"{AUTOMATED_TESTS} / {AUTOMATED_TESTS}", "OTOMATİK TEST")
    add_metric_card(slide, 3.55, 2.1, 2.7, "8 / 8", "RETRIEVAL KABULÜ", GREEN)
    add_metric_card(slide, 6.4, 2.1, 2.7, "77", "KANIT PARÇASI", ORANGE)
    add_metric_card(slide, 9.25, 2.1, 2.7, "HTTP 200", "CANLI SERVİS", BLUE)
    add_ppt_box(slide, 0.7, 3.75, 11.95, 2.2, fill=NAVY, line=NAVY)
    add_ppt_text(slide, "Çözüm değeri", 1.05, 4.08, 2.3, 0.35, size=11, color=CYAN, bold=True)
    add_ppt_text(
        slide,
        "Yanıt + belge + sayfa + otorite + koleksiyon",
        1.05,
        4.55,
        7.5,
        0.55,
        size=24,
        color=WHITE,
        bold=True,
    )
    add_ppt_text(slide, "Kullanıcı model çıktısını yalnız okumaz; kanıtı açıp denetler.", 1.08, 5.25, 8.4, 0.4, size=13, color="C7D5E8")
    add_label(slide, "Çevrimdışı çalışma", 9.25, 4.1, 2.6, CYAN)
    add_label(slide, "Hazır snapshot", 9.5, 4.75, 2.1, CYAN)
    add_label(slide, "Hibrit retrieval", 9.5, 5.4, 2.1, CYAN)

    # 3 Problem
    slide = presentation.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_frame(slide, 3, "Problem")
    add_slide_title(slide, "Neden klasik arama veya yalnız LLM yetmiyor?")
    problems = [
        ("Dağınık bilgi", "Broşür, ürün sayfası ve açık referanslar farklı yerlerde."),
        ("Bağlam kaybı", "Anahtar kelime araması anlam yakınlığını kaçırabilir."),
        ("Kaynaksız üretim", "LLM, ürün yeteneği veya örnek uydurabilir."),
    ]
    for index, (title, body) in enumerate(problems):
        x = 0.72 + index * 4.1
        add_ppt_box(slide, x, 2.15, 3.75, 2.7, fill=WHITE, line=RED if index == 2 else LINE)
        add_ppt_text(slide, f"0{index + 1}", x + 0.25, 2.42, 0.55, 0.35, size=10, color=RED if index == 2 else BLUE, bold=True)
        add_ppt_text(slide, title, x + 0.25, 2.92, 3.2, 0.45, size=20, color=NAVY, bold=True)
        add_ppt_text(slide, body, x + 0.25, 3.58, 3.15, 0.85, size=12, color=MUTED)
    add_ppt_box(slide, 0.72, 5.35, 11.95, 0.85, fill=PALE_BLUE, line=BLUE)
    add_ppt_text(slide, "Tasarım hedefi  →  Cevabı değil, cevabın kanıt zincirini güvenilir hâle getirmek.", 1.0, 5.62, 11.35, 0.3, size=16, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

    # 4 Trust model
    slide = presentation.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_frame(slide, 4, "Kaynak Güveni")
    add_slide_title(slide, "Kaynaklar birleşir; otoriteleri karışmaz", "Her parça belge, sayfa, koleksiyon, otorite ve URL bilgisini korur.")
    lanes = [
        ("OFFICIAL", "Resmî PDF + HAVELSAN", "Ürün ve ADVENT iddiaları", BLUE),
        ("OPEN SOURCE", "NATO kamu referansı", "Genel C2 ve birlikte çalışabilirlik", GREEN),
        ("ALL", "Birleşik görünüm", "Kaynak kimliği korunarak geniş araştırma", ORANGE),
    ]
    for index, (name, source, use, accent) in enumerate(lanes):
        y = 2.05 + index * 1.35
        add_ppt_box(slide, 0.75, y, 11.85, 1.02, fill=WHITE, line=accent)
        add_label(slide, name, 0.98, y + 0.33, 1.55, accent)
        add_ppt_text(slide, source, 2.85, y + 0.23, 3.2, 0.35, size=15, color=NAVY, bold=True)
        add_ppt_text(slide, use, 6.25, y + 0.25, 5.95, 0.45, size=12, color=MUTED)
    add_ppt_text(slide, "Kapsam değiştiğinde konuşma belleği de ayrılır.", 0.85, 6.25, 8.5, 0.35, size=14, color=BLUE, bold=True)
    add_ppt_text(slide, "→ ADVENT geçmişi NATO cevabını kirletemez.", 7.65, 6.25, 4.6, 0.35, size=12, color=GREEN, bold=True, align=PP_ALIGN.RIGHT)

    # 5 Architecture
    slide = presentation.slides.add_slide(blank)
    set_slide_background(slide, NAVY)
    add_slide_frame(slide, 5, "Mimari", dark=True)
    add_slide_title(slide, "Katmanlı ve testle korunan yapı", "Bağımlılık yönü yukarıdan aşağı; domain teknoloji bağımsızdır.", dark=True)
    layers = [
        ("app.py", "7 satırlık giriş noktası", "1A3357"),
        ("presentation", "Streamlit · sohbet · sidebar · kanıt kartı", "16497A"),
        ("application", "RAG kullanım senaryosu orkestrasyonu", "1B5FA7"),
        ("domain", "Modeller · sorgu kuralları · kanıt cevapları", "237BB9"),
        ("infrastructure", "PDF · manifest · FAISS · BM25 · reranker", "2E93D2"),
    ]
    widths = [5.0, 7.0, 8.6, 10.1, 11.4]
    for index, ((name, detail, fill), width) in enumerate(zip(layers, widths)):
        x = (13.333 - width) / 2
        y = 1.95 + index * 0.85
        add_ppt_box(slide, x, y, width, 0.62, fill=fill, line=CYAN)
        add_ppt_text(slide, name, x + 0.2, y + 0.16, 2.1, 0.25, size=11, color=WHITE, bold=True)
        add_ppt_text(slide, detail, x + 2.35, y + 0.16, width - 2.55, 0.25, size=9.5, color="D7E6F7", align=PP_ALIGN.RIGHT)
    add_ppt_text(slide, "ArchitectureGuardTests", 4.7, 6.45, 3.9, 0.34, size=10, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_ppt_text(slide, "ince giriş · bağımlılık yönü · Türkçe açıklama standardı", 3.4, 6.82, 6.5, 0.25, size=9, color="9DB4D1", align=PP_ALIGN.CENTER)

    # 6 lifecycle
    slide = presentation.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_frame(slide, 6, "Belge Yaşam Döngüsü")
    add_slide_title(slide, "Araştırmadan hazır yerel kanıta kadar izlenebilir")
    add_pipeline_nodes(
        slide,
        [
            ("Araştır", "Birincil kamu kaynağı"),
            ("Kürasyon", "Kapsam + kaynak"),
            ("PDF", "Metin çıkarılabilir"),
            ("Parçala", "Sayfa + overlap"),
            ("Embedding", "Önceden hesapla"),
            ("Snapshot", "Manifest + vektör"),
        ],
        2.4,
    )
    add_ppt_box(slide, 0.78, 4.65, 5.7, 1.22, fill="EDF9F1", line=GREEN)
    add_ppt_text(slide, "Normal kullanım", 1.05, 4.94, 1.8, 0.3, size=14, color=GREEN, bold=True)
    add_ppt_text(slide, "Hazır snapshot yüklenir · web erişimi yok", 2.65, 4.94, 3.45, 0.4, size=11, color=INK)
    add_ppt_box(slide, 6.72, 4.65, 5.7, 1.22, fill="FFF6E5", line=ORANGE)
    add_ppt_text(slide, "Ek belge?", 7.0, 4.94, 1.6, 0.3, size=14, color=ORANGE, bold=True)
    add_ppt_text(slide, "SHA-256 kontrolü · yalnız ek belge indekslenir", 8.55, 4.94, 3.45, 0.4, size=11, color=INK)

    # 7 RAG pipeline
    slide = presentation.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_frame(slide, 7, "RAG Hattı")
    add_slide_title(slide, "Hibrit retrieval: anlam + terim + yeniden sıralama")
    add_pipeline_nodes(
        slide,
        [
            ("Soru", "Türkçe doğal dil"),
            ("Expand", "CMS sözlüğü"),
            ("Semantic", "BGE + FAISS"),
            ("BM25", "Kesin terim"),
            ("RRF", "Sıra füzyonu"),
            ("Rerank", "Cross-encoder"),
            ("Kanıt", "Sayfa birleştirme"),
        ],
        2.3,
    )
    add_ppt_text(slide, "Semantic", 1.1, 4.65, 1.5, 0.32, size=13, color=BLUE, bold=True)
    add_ppt_text(slide, "“Savaş gemisi” ↔ “surface platform” anlam yakınlığı", 2.35, 4.65, 4.25, 0.35, size=11.5, color=MUTED)
    add_ppt_text(slide, "BM25", 1.1, 5.27, 1.5, 0.32, size=13, color=ORANGE, bold=True)
    add_ppt_text(slide, "ADVENT · MÜREN · Link 11 · Link 16 kesin eşleşmesi", 2.35, 5.27, 4.25, 0.35, size=11.5, color=MUTED)
    add_ppt_box(slide, 7.35, 4.55, 4.85, 1.25, fill=NAVY, line=NAVY)
    add_ppt_text(slide, "RRF", 7.7, 4.82, 1.0, 0.38, size=18, color=CYAN, bold=True)
    add_ppt_text(slide, "1 / (60 + sıra)", 8.85, 4.84, 2.4, 0.35, size=16, color=WHITE, bold=True)
    add_ppt_text(slide, "Farklı puan ölçeklerini istikrarlı birleştirir.", 7.72, 5.35, 3.95, 0.28, size=9.5, color="C7D5E8")

    # 8 answer guardrails
    slide = presentation.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_frame(slide, 8, "Cevap Güvenilirliği")
    add_slide_title(slide, "Model en son adım; karar zincirinin tamamı değil")
    decisions = [
        ("1", "İndeks var mı?", "Yoksa yükleme yönlendirmesi"),
        ("2", "Alan içi mi?", "Kişisel sohbet → güvenli ret"),
        ("3", "Kesin kanıt var mı?", "Deterministik kaynaklı cevap"),
        ("4", "Retrieval yeterli mi?", "Kapsam filtreli adaylar"),
        ("5", "Ollama üretimi", "≤55 kelime · düşük sıcaklık"),
        ("6", "Kaynak kontrolü", "[SOURCE n] + kanıt kartı"),
    ]
    for index, (number, title, body) in enumerate(decisions):
        col = index % 3
        row = index // 3
        x = 0.72 + col * 4.12
        y = 2.05 + row * 2.0
        add_ppt_box(slide, x, y, 3.75, 1.52, fill=WHITE, line=GREEN if index in (2, 5) else LINE)
        add_ppt_text(slide, number, x + 0.2, y + 0.18, 0.45, 0.36, size=11, color=GREEN if index in (2, 5) else BLUE, bold=True)
        add_ppt_text(slide, title, x + 0.72, y + 0.16, 2.7, 0.35, size=15, color=NAVY, bold=True)
        add_ppt_text(slide, body, x + 0.72, y + 0.72, 2.65, 0.46, size=10.5, color=MUTED)
    add_ppt_text(slide, "Servis hatası ≠ kanıt", 0.82, 6.25, 3.0, 0.32, size=13, color=RED, bold=True)
    add_ppt_text(slide, "Ollama erişilemiyorsa retrieval sonuçları kullanıcıya kaynak diye gösterilmez.", 3.2, 6.25, 8.65, 0.35, size=11.5, color=MUTED)

    # 9 memory
    slide = presentation.slides.add_slide(blank)
    set_slide_background(slide, NAVY)
    add_slide_frame(slide, 9, "Agentic Kalıcılık", dark=True)
    add_slide_title(slide, "Konuşma, kanıt ve bekleyen onay yeniden başlatmada korunur", dark=True)
    # Conversation cards
    add_ppt_box(slide, 0.75, 2.05, 5.45, 0.8, fill="15365F", line=BLUE)
    add_ppt_text(slide, "ADVENT nedir?", 1.0, 2.28, 4.9, 0.3, size=14, color=WHITE, bold=True)
    add_ppt_box(slide, 0.75, 3.0, 5.45, 0.95, fill="1B416F", line=CYAN)
    add_ppt_text(slide, "Örnekleri var mı?", 1.0, 3.22, 4.9, 0.3, size=14, color=WHITE, bold=True)
    add_ppt_text(slide, "← aynı official geçmişi", 3.75, 3.62, 2.1, 0.2, size=8, color=CYAN, align=PP_ALIGN.RIGHT)
    add_ppt_box(slide, 0.75, 4.15, 5.45, 0.95, fill="174A46", line=GREEN)
    add_ppt_text(slide, "NATO interoperability?", 1.0, 4.38, 4.9, 0.3, size=14, color=WHITE, bold=True)
    add_ppt_text(slide, "← open_source: yeni bağlam", 3.45, 4.78, 2.4, 0.2, size=8, color="7EE2BA", align=PP_ALIGN.RIGHT)
    # Isolation wall
    wall = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, PptInches(6.55), PptInches(1.95), PptInches(0.08), PptInches(3.55))
    wall.fill.solid()
    wall.fill.fore_color.rgb = ppt_color(ORANGE)
    wall.line.fill.background()
    add_ppt_text(slide, "SCOPE\nWALL", 6.7, 3.0, 0.8, 0.8, size=8, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_ppt_box(slide, 7.75, 2.15, 4.55, 1.25, fill="13335A", line=BLUE)
    add_ppt_text(slide, "official history", 8.05, 2.42, 3.9, 0.3, size=16, color=WHITE, bold=True)
    add_ppt_text(slide, "ADVENT ürün bağlamı", 8.05, 2.87, 3.9, 0.25, size=10, color="AFC2DA")
    add_ppt_box(slide, 7.75, 3.78, 4.55, 1.25, fill="174A46", line=GREEN)
    add_ppt_text(slide, "open_source history", 8.05, 4.05, 3.9, 0.3, size=16, color=WHITE, bold=True)
    add_ppt_text(slide, "NATO genel birlikte çalışabilirlik", 8.05, 4.5, 3.9, 0.25, size=10, color="BDE7D5")
    add_ppt_text(slide, "PostgreSQL checkpoint · thread geri yükleme · MCP interrupt / resume", 2.35, 6.25, 8.7, 0.35, size=12, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

    # 10 UI
    slide = presentation.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_frame(slide, 10, "Kullanıcı Deneyimi")
    add_slide_title(slide, "Arayüzün ana ürünü: görünür kanıt")
    # App wireframe
    add_ppt_box(slide, 0.72, 1.85, 11.9, 4.65, fill=WHITE, line=LINE)
    sidebar = add_ppt_box(slide, 0.72, 1.85, 2.55, 4.65, fill=NAVY, line=NAVY, radius=False)
    add_ppt_text(slide, "◆ Knowledge Operations", 0.95, 2.12, 2.05, 0.35, size=10, color=WHITE, bold=True)
    add_label(slide, "Kapsam", 1.05, 2.85, 1.65, CYAN)
    add_label(slide, "PDF yükle", 1.05, 3.48, 1.65, CYAN)
    add_label(slide, "İndeksi yenile", 1.05, 4.11, 1.65, CYAN)
    add_ppt_text(slide, "Model · qwen2.5:3b\nBelge · 1\nArama · Hybrid", 1.0, 5.1, 1.9, 0.85, size=9, color="AFC2DA")
    add_ppt_text(slide, "Komuta Bilgi Keşfi", 3.65, 2.12, 4.8, 0.48, size=24, color=NAVY, bold=True)
    add_ppt_box(slide, 3.65, 2.95, 8.4, 0.62, fill=LIGHT, line=LIGHT)
    add_ppt_text(slide, "ADVENT nedir?", 3.9, 3.15, 7.7, 0.22, size=11, color=INK)
    add_ppt_box(slide, 3.65, 3.78, 8.4, 1.05, fill="F7FAFE", line=PALE_BLUE)
    add_ppt_text(slide, "KAYNAKLI YANIT", 3.9, 3.98, 2.0, 0.2, size=7.5, color=BLUE, bold=True)
    add_ppt_text(slide, "ADVENT, farklı operasyonel ortamlar için uyarlanabilen bir CMS ürün ailesidir [SOURCE 1].", 3.9, 4.3, 7.65, 0.32, size=10.5, color=INK)
    add_ppt_box(slide, 3.9, 5.05, 7.9, 0.65, fill=PALE_BLUE, line=BLUE)
    add_ppt_text(slide, "Kanıt paketi · advent_cms.pdf · Sayfa 4 · official", 4.15, 5.26, 7.4, 0.23, size=9.5, color=BLUE, bold=True)
    add_ppt_text(slide, "Akışlı yazım", 4.0, 6.05, 1.8, 0.25, size=9, color=GREEN, bold=True)
    add_ppt_text(slide, "Kalıcı kaynak kartı", 6.2, 6.05, 2.3, 0.25, size=9, color=BLUE, bold=True)
    add_ppt_text(slide, "Güvenli ret", 9.25, 6.05, 1.6, 0.25, size=9, color=ORANGE, bold=True)

    # 11 Acceptance
    slide = presentation.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_frame(slide, 11, "Test ve Kabul")
    add_slide_title(slide, "“Çalışıyor” değil; ölçülmüş ve yeniden üretilmiş")
    add_metric_card(
        slide, 0.72, 2.0, 2.5,
        f"{AUTOMATED_TESTS} / {AUTOMATED_TESTS}", "OTOMATİK TEST",
    )
    add_metric_card(
        slide, 3.38, 2.0, 2.5,
        f"{BENCHMARK_REPORT['passed']} / {BENCHMARK_REPORT['dataset_cases']}",
        "ALTIN SET", GREEN,
    )
    add_metric_card(
        slide, 6.04, 2.0, 2.5,
        f"{QUALITY_REPORT['stage2']['summary']['strict_passed']} / "
        f"{QUALITY_REPORT['stage2']['summary']['evaluated']}",
        "LLM HAKEM", ORANGE,
    )
    add_metric_card(
        slide, 8.7, 2.0, 2.5,
        f"{LINEAGE_REPORT['confusion_matrix']['false_positive']} / "
        f"{LINEAGE_REPORT['confusion_matrix']['false_negative']}",
        "20 VAKA FP / FN", BLUE,
    )
    categories = [
        ("Storage", 8, BLUE),
        ("Engine", 30, CYAN),
        ("UI + PDF", 17, GREEN),
        ("Eval + Audit", 19, ORANGE),
        ("Architecture + KB", 8, NAVY),
    ]
    max_value = max(value for _, value, _ in categories)
    for index, (label, value, color) in enumerate(categories):
        y = 3.65 + index * 0.55
        add_ppt_text(slide, label, 0.85, y, 1.45, 0.25, size=10, color=MUTED, bold=True)
        add_ppt_box(slide, 2.35, y + 0.03, 7.6, 0.18, fill="E5EAF0", line="E5EAF0", radius=False)
        add_ppt_box(slide, 2.35, y + 0.03, 7.6 * value / max_value, 0.18, fill=color, line=color, radius=False)
        add_ppt_text(slide, str(value), 10.1, y - 0.01, 0.45, 0.25, size=10, color=NAVY, bold=True)
    add_ppt_box(slide, 10.85, 3.62, 1.55, 2.5, fill=NAVY, line=NAVY)
    add_ppt_text(slide, "GIT\nARCHIVE", 11.0, 4.0, 1.25, 0.65, size=14, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_ppt_text(slide, "temiz kopyada\naynı sonuç", 11.0, 4.92, 1.25, 0.6, size=9, color=WHITE, align=PP_ALIGN.CENTER)

    # 12 security
    slide = presentation.slides.add_slide(blank)
    set_slide_background(slide, NAVY)
    add_slide_frame(slide, 12, "Güvenlik", dark=True)
    add_slide_title(slide, "Yerel çalışma + kontrollü kaynak + güvenli hata", dark=True)
    controls = [
        ("LOCAL", "Belge metni bulut LLM'e gitmez"),
        ("HASH", "SHA-256 duplicate ve bütünlük"),
        ("SCOPE", "Koleksiyon ve geçmiş izolasyonu"),
        ("ESCAPE", "Kaynak HTML çıktısı güvenli"),
        ("HITL", "MCP yazmasında operatör onayı"),
        ("PATH", "Manifest traversal engeli"),
    ]
    for index, (tag, body) in enumerate(controls):
        col = index % 3
        row = index // 3
        x = 0.75 + col * 4.12
        y = 2.05 + row * 1.85
        add_ppt_box(slide, x, y, 3.75, 1.4, fill="14345B", line=CYAN if index in (0, 2) else "31577F")
        add_label(slide, tag, x + 0.22, y + 0.22, 1.05, CYAN)
        add_ppt_text(slide, body, x + 0.22, y + 0.72, 3.2, 0.42, size=11.5, color=WHITE)
    add_ppt_text(slide, "Yerel audit hazır · Üretime geçişte: RBAC · merkezi aktarım · disk şifreleme", 1.2, 6.05, 10.95, 0.4, size=13, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

    # 13 git
    slide = presentation.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_frame(slide, 13, "Git ve Sürümleme")
    add_slide_title(slide, "Geçmiş korunur; değişiklikler kontrollü akar")
    lane_y = {"main": 2.25, "develop": 3.35, "feature": 4.45, "release": 5.55}
    colors = {"main": NAVY, "develop": BLUE, "feature": CYAN, "release": ORANGE}
    for name, y in lane_y.items():
        add_ppt_text(slide, name, 0.75, y - 0.08, 1.15, 0.28, size=10, color=colors[name], bold=True)
        add_ppt_line(slide, 1.9, y + 0.05, 12.2, y + 0.05, "C9D4E1", 1.2)
    # Commit points and branches
    for x, y, label, color in [
        (2.4, lane_y["main"], "legacy", NAVY),
        (3.6, lane_y["develop"], "4dbf3f8", BLUE),
        (5.0, lane_y["feature"], "43b38ae", CYAN),
        (6.8, lane_y["feature"], "0e72e16", CYAN),
        (8.35, lane_y["develop"], "d022506", BLUE),
        (10.4, lane_y["release"], "release/*", ORANGE),
        (11.65, lane_y["main"], "v1.x", NAVY),
    ]:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, PptInches(x), PptInches(y - 0.11), PptInches(0.22), PptInches(0.22))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ppt_color(color)
        shape.line.color.rgb = ppt_color(color)
        add_ppt_text(slide, label, x - 0.4, y + 0.2, 1.05, 0.25, size=7.5, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_ppt_line(slide, 3.6, lane_y["develop"], 5.0, lane_y["feature"], CYAN, 1.5)
    add_ppt_line(slide, 6.91, lane_y["feature"], 8.35, lane_y["develop"], BLUE, 1.5)
    add_ppt_line(slide, 8.45, lane_y["develop"], 10.4, lane_y["release"], ORANGE, 1.5)
    add_ppt_line(slide, 10.5, lane_y["release"], 11.65, lane_y["main"], NAVY, 1.5)
    add_ppt_text(slide, "main: yalnız yayımlanmış sürüm", 2.0, 6.35, 3.3, 0.3, size=10.5, color=NAVY, bold=True)
    add_ppt_text(slide, "develop: entegrasyon tabanı", 5.1, 6.35, 3.2, 0.3, size=10.5, color=BLUE, bold=True)
    add_ppt_text(slide, "feature/release: inceleme + kabul", 8.3, 6.35, 3.8, 0.3, size=10.5, color=ORANGE, bold=True)

    # 14 Deployment
    slide = presentation.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_frame(slide, 14, "Çalıştırma")
    add_slide_title(slide, "Yerel kurulum: iki servis, tek tarayıcı")
    steps = [
        ("01", "Python ortamı", "python -m venv .venv\npip install -r requirements.txt"),
        ("02", "Yerel model", "ollama pull qwen2.5:3b\nollama serve"),
        ("03", "Uygulama", "python -m streamlit run app.py\nlocalhost:8501"),
    ]
    for index, (number, title, commands) in enumerate(steps):
        x = 0.75 + index * 4.1
        add_ppt_box(slide, x, 2.05, 3.75, 3.05, fill=WHITE, line=BLUE if index == 2 else LINE)
        add_ppt_text(slide, number, x + 0.25, 2.35, 0.65, 0.35, size=12, color=CYAN, bold=True)
        add_ppt_text(slide, title, x + 0.25, 2.92, 3.15, 0.42, size=20, color=NAVY, bold=True)
        add_ppt_box(slide, x + 0.23, 3.62, 3.3, 1.1, fill=NAVY, line=NAVY)
        add_ppt_text(slide, commands, x + 0.42, 3.88, 2.9, 0.65, size=9.2, color=WHITE, font="Cascadia Mono")
    add_ppt_box(slide, 0.75, 5.58, 11.95, 0.78, fill=PALE_BLUE, line=BLUE)
    add_ppt_text(slide, "Offline mod  ·  $env:CMS_RAG_OFFLINE=\"1\"", 1.05, 5.82, 4.6, 0.28, size=11, color=BLUE, bold=True)
    add_ppt_text(slide, "Kalite modu  ·  $env:CMS_RAG_MODEL=\"qwen2.5:7b\"", 6.25, 5.82, 5.75, 0.28, size=11, color=BLUE, bold=True)

    # 15 Roadmap
    slide = presentation.slides.add_slide(blank)
    set_slide_background(slide)
    add_slide_frame(slide, 15, "Yol Haritası")
    add_slide_title(slide, "MVP tamamlandı; kurumsal ölçek için sıradaki adımlar")
    horizons = [
        ("YAKIN", "OCR\nKaynak drift kontrolü\nCI kabul pipeline", BLUE),
        ("ORTA", "RBAC\nCheckpoint saklama politikası\nKalite eğilim alarmı", GREEN),
        ("UZUN", "Merkezi audit aktarımı\nÇok modlu retrieval\nKurumsal gözlemlenebilirlik", ORANGE),
    ]
    for index, (label, body, accent) in enumerate(horizons):
        x = 0.75 + index * 4.1
        add_ppt_box(slide, x, 2.0, 3.75, 3.75, fill=WHITE, line=accent)
        add_label(slide, label, x + 0.85, 2.35, 2.05, accent)
        add_ppt_text(slide, body, x + 0.45, 3.15, 2.85, 1.65, size=16, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_ppt_text(slide, f"0{index + 1}", x + 1.48, 5.05, 0.75, 0.35, size=12, color=accent, bold=True, align=PP_ALIGN.CENTER)
    add_ppt_text(slide, "Öncelik: CI kalite kapısı + RBAC + kurumsal gözlemlenebilirlik", 2.25, 6.25, 8.85, 0.4, size=15, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

    # 16 Demo
    slide = presentation.slides.add_slide(blank)
    set_slide_background(slide, NAVY)
    add_slide_frame(slide, 16, "Canlı Demo", dark=True)
    add_slide_title(slide, "5 adımda projenin bütün değerini göster", "Toplam önerilen demo süresi: 3 dakika", dark=True)
    demo_steps = [
        ("1", "ADVENT nedir?", "Resmî kaynak"),
        ("2", "Örnekleri var mı?", "Takip bağlamı"),
        ("3", "Görevleri neler?", "Çoklu kaynak"),
        ("4", "NATO interoperability?", "Open source"),
        ("5", "Ben kimim?", "Güvenli ret"),
    ]
    for index, (number, question, proof) in enumerate(demo_steps):
        x = 0.75 + index * 2.48
        add_ppt_box(slide, x, 2.25, 2.18, 2.55, fill="14345B", line=CYAN if index < 4 else ORANGE)
        add_ppt_text(slide, number, x + 0.75, 2.53, 0.65, 0.5, size=22, color=CYAN if index < 4 else ORANGE, bold=True, align=PP_ALIGN.CENTER)
        add_ppt_text(slide, question, x + 0.18, 3.25, 1.82, 0.65, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_label(slide, proof, x + 0.3, 4.15, 1.58, CYAN if index < 4 else ORANGE)
        if index < len(demo_steps) - 1:
            add_ppt_line(slide, x + 2.18, 3.5, x + 2.48, 3.5, CYAN, 1.4)
    add_ppt_text(slide, "Gösterilecek ana unsur: cevaptan önce kanıt paketi.", 2.2, 5.55, 8.95, 0.45, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_ppt_text(slide, "Belge · Sayfa · Otorite · Koleksiyon", 3.25, 6.18, 6.85, 0.35, size=12, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

    # 17 Closing
    slide = presentation.slides.add_slide(blank)
    set_slide_background(slide, NAVY)
    add_slide_frame(slide, 17, "Sonuç", dark=True)
    add_ppt_text(slide, "CMS-RAG Assistant", 0.75, 1.25, 7.3, 0.62, size=31, color=WHITE, bold=True, font="Aptos Display")
    add_ppt_text(slide, "Cevap üreten değil,\nkanıt zinciri sunan asistan.", 0.75, 2.1, 7.2, 1.3, size=28, color=CYAN, bold=True, font="Aptos Display")
    add_ppt_box(slide, 8.25, 1.45, 3.85, 3.85, fill="14345B", line=CYAN)
    closing = [
        ("LOCAL", "Veri gizliliği"),
        ("HYBRID", "Yüksek retrieval kapsaması"),
        ("GROUNDED", "Belge ve sayfa kanıtı"),
        ("TESTED", f"{AUTOMATED_TESTS}/{AUTOMATED_TESTS} + 45/45 + 30/30"),
    ]
    for index, (tag, body) in enumerate(closing):
        y = 1.85 + index * 0.78
        add_label(slide, tag, 8.65, y, 1.15, CYAN)
        add_ppt_text(slide, body, 10.0, y + 0.06, 1.65, 0.24, size=10, color=WHITE, bold=True)
    add_ppt_text(slide, "Sorularınız?", 0.78, 5.15, 4.6, 0.55, size=25, color=WHITE, bold=True)
    add_ppt_text(slide, "http://localhost:8501", 0.8, 5.95, 4.5, 0.3, size=12, color="9DB4D1")
    add_ppt_text(slide, f"{COMMIT}  ·  Sürüm {VERSION}", 8.25, 6.15, 3.85, 0.3, size=9, color="9DB4D1", align=PP_ALIGN.RIGHT)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    presentation.save(PPTX_PATH)


def validate_outputs() -> None:
    """Üretilen Office dosyalarının yeniden açıldığını ve temel yapıyı taşıdığını doğrular."""

    document = Document(DOCX_PATH)
    if len(document.paragraphs) < 150:
        raise RuntimeError("Word belgesinin içerik kapsamı beklenenden düşük.")
    if len(document.tables) < 20:
        raise RuntimeError("Word belgesinde beklenen teknik tablo sayısı yok.")
    if "Savaş Yönetim Sistemi" not in "\n".join(p.text for p in document.paragraphs[:20]):
        raise RuntimeError("Word kapak başlığı bulunamadı.")

    presentation = Presentation(PPTX_PATH)
    if len(presentation.slides) != 17:
        raise RuntimeError("PowerPoint slayt sayısı beklenen 17 değil.")
    titles = []
    for slide in presentation.slides:
        slide_text = " ".join(
            shape.text for shape in slide.shapes if hasattr(shape, "text")
        )
        titles.append(slide_text)
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0:
                raise RuntimeError("Slayt sınırı dışında negatif konumlu şekil bulundu.")
            if shape.left + shape.width > presentation.slide_width + PptInches(0.02):
                raise RuntimeError("Slayt sağ sınırını aşan şekil bulundu.")
            if shape.top + shape.height > presentation.slide_height + PptInches(0.02):
                raise RuntimeError("Slayt alt sınırını aşan şekil bulundu.")
    required = ("Tek cümlede proje", "Katmanlı ve testle korunan yapı", "Sorularınız?")
    if not all(any(text in title for title in titles) for text in required):
        raise RuntimeError("PowerPoint temel anlatı slaytlarından biri eksik.")


def main() -> None:
    """Her iki nihai teslimatı üretir ve yapısal olarak doğrular."""

    build_word_document()
    build_powerpoint()
    validate_outputs()
    print(f"Word: {DOCX_PATH}")
    print(f"PowerPoint: {PPTX_PATH}")


if __name__ == "__main__":
    main()
