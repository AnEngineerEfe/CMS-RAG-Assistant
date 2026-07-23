"""Generate the CMS-RAG project presentation as a polished .pptx deck."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as Shape
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


OUT = Path("docs/CMS-RAG-Proje-Sunumu.pptx")
W, H = Inches(13.333), Inches(7.5)
NAVY = "0B1F33"
NAVY_2 = "123047"
TEAL = "00A6A6"
CYAN = "40D5D5"
WHITE = "F7FBFC"
MUTED = "A9BBC5"
PANEL = "173B53"
GREEN = "47D18C"
AMBER = "F8BE4A"


def rgb(value):
    return RGBColor.from_string(value)


def add_shape(slide, kind, x, y, w, h, fill, line=None, radius=False):
    shape = slide.shapes.add_shape(Shape.ROUNDED_RECTANGLE if radius else kind, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    return shape


def textbox(slide, text, x, y, w, h, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Aptos", valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def line(slide, x1, y1, x2, y2, color=TEAL, width=2):
    item = slide.shapes.add_connector(1, x1, y1, x2, y2)
    item.line.color.rgb = rgb(color)
    item.line.width = Pt(width)
    return item


def background(slide, number, title=None, eyebrow="CMS-RAG ASSISTANT"):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(NAVY)
    add_shape(slide, Shape.RECTANGLE, 0, 0, W, Inches(0.12), TEAL)
    textbox(slide, eyebrow, Inches(0.62), Inches(0.32), Inches(4), Inches(0.24), 9, CYAN, True)
    if title:
        textbox(slide, title, Inches(0.62), Inches(0.64), Inches(11.9), Inches(0.58), 27, WHITE, True)
    textbox(slide, f"{number:02d}", Inches(12.1), Inches(7.02), Inches(0.6), Inches(0.25), 9, MUTED, True, PP_ALIGN.RIGHT)
    textbox(slide, "Yerel • Kaynak Gösteren • Denetlenebilir", Inches(0.62), Inches(7.02), Inches(5), Inches(0.25), 9, MUTED)


def card(slide, title, body, x, y, w, h, accent=TEAL, icon=None):
    add_shape(slide, Shape.ROUNDED_RECTANGLE, x, y, w, h, PANEL, PANEL, True)
    add_shape(slide, Shape.RECTANGLE, x, y, Inches(0.07), h, accent)
    if icon:
        textbox(slide, icon, x + Inches(0.26), y + Inches(0.24), Inches(0.48), Inches(0.45), 21, accent, True, PP_ALIGN.CENTER)
        tx = x + Inches(0.82)
    else:
        tx = x + Inches(0.28)
    textbox(slide, title, tx, y + Inches(0.24), w - (tx - x) - Inches(0.2), Inches(0.36), 15, WHITE, True)
    textbox(slide, body, x + Inches(0.28), y + Inches(0.77), w - Inches(0.56), h - Inches(0.95), 11.5, MUTED)


def chip(slide, text, x, y, w, fill=TEAL):
    shape = add_shape(slide, Shape.ROUNDED_RECTANGLE, x, y, w, Inches(0.34), fill, fill, True)
    textbox(slide, text, x, y + Inches(0.05), w, Inches(0.2), 9, NAVY, True, PP_ALIGN.CENTER)
    return shape


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    blank = prs.slide_layouts[6]

    # 1 — Cover
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(NAVY)
    add_shape(slide, Shape.RECTANGLE, 0, 0, W, Inches(0.14), TEAL)
    # Abstract maritime/grid visual
    for i in range(7):
        line(slide, Inches(7.3 + i * 0.72), Inches(1.1), Inches(12.9), Inches(1.1 + i * 0.72), PANEL, 1)
        line(slide, Inches(7.3), Inches(1.1 + i * 0.72), Inches(7.3 + i * 0.72), Inches(5.42), PANEL, 1)
    for x, y, r in [(9.7, 2.2, 0.75), (11.05, 3.25, 0.48), (9.1, 4.52, 0.42), (11.8, 5.1, 0.28)]:
        shape = add_shape(slide, Shape.OVAL, Inches(x-r/2), Inches(y-r/2), Inches(r), Inches(r), TEAL, TEAL)
        shape.fill.transparency = 25
        line(slide, Inches(9.7), Inches(2.2), Inches(x), Inches(y), CYAN, 1)
    textbox(slide, "CMS-RAG\nASİSTANI", Inches(0.75), Inches(1.45), Inches(6.2), Inches(1.65), 42, WHITE, True, font="Aptos Display")
    textbox(slide, "Savaş Yönetim Sistemi bilgi erişiminde\nyerel, kaynak gösteren ve denetlenebilir yapay zekâ", Inches(0.78), Inches(3.38), Inches(5.8), Inches(0.82), 17, MUTED)
    chip(slide, "ADVENT CMS", Inches(0.78), Inches(4.64), Inches(1.35), CYAN)
    chip(slide, "HYBRID RAG", Inches(2.28), Inches(4.64), Inches(1.45), TEAL)
    chip(slide, "ON-PREMISE", Inches(3.88), Inches(4.64), Inches(1.6), GREEN)
    textbox(slide, "Proje Sunumu • Temmuz 2026", Inches(0.78), Inches(6.65), Inches(4), Inches(0.25), 10, MUTED)

    # 2 — Problem
    slide = prs.slides.add_slide(blank); background(slide, 2, "Problem: Bilgi Çok, Güvenilir Erişim Zor")
    textbox(slide, "CMS dokümanları kritik, dağınık ve yüksek doğruluk gerektiriyor.", Inches(0.62), Inches(1.35), Inches(11), Inches(0.35), 15, MUTED)
    card(slide, "Dağınık Bilgi", "Broşürler, teknik kılavuzlar, standartlar ve kamu referansları farklı yerlerde bulunur.", Inches(0.68), Inches(2.0), Inches(3.8), Inches(2.1), AMBER, "01")
    card(slide, "Kaynak Riski", "Genel amaçlı LLM'ler, doğrulanmamış veya kaynaksız içerik üretebilir.", Inches(4.77), Inches(2.0), Inches(3.8), Inches(2.1), "E95C65", "02")
    card(slide, "Gizlilik Riski", "Hassas/kurumsal verinin haricî servislerle paylaşılması kabul edilemez.", Inches(8.86), Inches(2.0), Inches(3.8), Inches(2.1), "A26EEA", "03")
    add_shape(slide, Shape.ROUNDED_RECTANGLE, Inches(1.5), Inches(4.85), Inches(10.3), Inches(0.9), "102A41", "102A41", True)
    textbox(slide, "Hedef: Doğru bilgiye, doğru kaynakla ve doğru güven sınırı içinde erişim.", Inches(1.7), Inches(5.1), Inches(9.9), Inches(0.35), 19, CYAN, True, PP_ALIGN.CENTER)

    # 3 — Value
    slide = prs.slides.add_slide(blank); background(slide, 3, "Çözüm: Kanıt Tabanlı Yerel Asistan")
    steps = [("1", "TOPLA", "Onaylı kaynaklar"), ("2", "BUL", "Hibrit arama"), ("3", "SEÇ", "Reranking"), ("4", "YANITLA", "Yerel LLM")]
    for i, (num, title, body) in enumerate(steps):
        x = Inches(0.72 + i * 3.12)
        add_shape(slide, Shape.OVAL, x + Inches(0.87), Inches(1.85), Inches(1.05), Inches(1.05), TEAL if i < 3 else GREEN, TEAL)
        textbox(slide, num, x + Inches(0.87), Inches(2.11), Inches(1.05), Inches(0.35), 20, NAVY, True, PP_ALIGN.CENTER)
        textbox(slide, title, x, Inches(3.2), Inches(2.8), Inches(0.35), 16, WHITE, True, PP_ALIGN.CENTER)
        textbox(slide, body, x, Inches(3.65), Inches(2.8), Inches(0.35), 11, MUTED, False, PP_ALIGN.CENTER)
        if i < 3:
            line(slide, x + Inches(2.1), Inches(2.37), x + Inches(3.03), Inches(2.37), CYAN, 2)
    card(slide, "Çıktı", "Teknik, kısa ve kaynakları görünür yanıt. Kullanıcı belge, sayfa ve alaka skorunu doğrudan görür.", Inches(2.25), Inches(5.0), Inches(8.85), Inches(1.15), TEAL, "✓")

    # 4 Architecture
    slide = prs.slides.add_slide(blank); background(slide, 4, "Kurumsal Mimari: İki Koleksiyon, Tek Orkestrasyon")
    card(slide, "HAVELSAN", "Resmî ADVENT broşürleri\nve yetkili dokümanlar", Inches(0.72), Inches(1.65), Inches(2.6), Inches(1.3), CYAN, "A")
    card(slide, "Açık / Kamu", "NATO, NISP ve\ndoğrulanmış referanslar", Inches(0.72), Inches(3.4), Inches(2.6), Inches(1.3), GREEN, "B")
    line(slide, Inches(3.38), Inches(2.3), Inches(4.32), Inches(2.95), CYAN, 2)
    line(slide, Inches(3.38), Inches(4.05), Inches(4.32), Inches(3.3), GREEN, 2)
    add_shape(slide, Shape.ROUNDED_RECTANGLE, Inches(4.35), Inches(2.25), Inches(2.5), Inches(1.7), "102F47", TEAL, True)
    textbox(slide, "CMS\nKNOWLEDGE BASE", Inches(4.55), Inches(2.66), Inches(2.1), Inches(0.68), 17, WHITE, True, PP_ALIGN.CENTER)
    line(slide, Inches(6.92), Inches(3.1), Inches(7.78), Inches(3.1), TEAL, 2)
    card(slide, "Hybrid Retrieval", "FAISS + BM25\nAday bulma", Inches(7.83), Inches(1.6), Inches(2.1), Inches(1.3), TEAL, "R")
    card(slide, "Reranker", "Cross-encoder\nEn iyi bağlam", Inches(7.83), Inches(3.38), Inches(2.1), Inches(1.3), AMBER, "R")
    line(slide, Inches(10.0), Inches(2.25), Inches(10.66), Inches(2.93), CYAN, 2)
    line(slide, Inches(10.0), Inches(4.03), Inches(10.66), Inches(3.32), AMBER, 2)
    card(slide, "Ollama", "Qwen 2.5:3b\nYerel yanıt", Inches(10.7), Inches(2.25), Inches(2.0), Inches(1.7), GREEN, "L")
    textbox(slide, "Resmî bilgi, genel referanslardan hiçbir zaman sessizce karışmaz.", Inches(1.15), Inches(5.65), Inches(11.2), Inches(0.4), 17, CYAN, True, PP_ALIGN.CENTER)

    # 5 Data and governance
    slide = prs.slides.add_slide(blank); background(slide, 5, "Veri Yönetişimi ve Güven Sınırı")
    textbox(slide, "Koleksiyon ayrımı; erişim kalitesinin yanında denetlenebilirlik için tasarım kararıdır.", Inches(0.62), Inches(1.35), Inches(11.8), Inches(0.35), 14, MUTED)
    card(slide, "Resmî Koleksiyon", "• collection: havelsan\n• authority: official\n• Ürün iddiaları için birincil kaynak", Inches(0.72), Inches(2.0), Inches(3.85), Inches(2.15), CYAN, "01")
    card(slide, "Kamu Referansları", "• collection: open_source\n• NATO/NISP bağlamı\n• Standart ve kavram açıklamaları", Inches(4.75), Inches(2.0), Inches(3.85), Inches(2.15), GREEN, "02")
    card(slide, "Kurum İçi Belgeler", "• collection: uploaded\n• Onay ve sınıflandırma gerekir\n• İzinli yerel ortam", Inches(8.78), Inches(2.0), Inches(3.85), Inches(2.15), AMBER, "03")
    add_shape(slide, Shape.ROUNDED_RECTANGLE, Inches(1.35), Inches(4.95), Inches(10.65), Inches(0.9), "102A41", "102A41", True)
    textbox(slide, "Her parça: belge adı • kaynak yolu • koleksiyon • otorite • sayfa", Inches(1.55), Inches(5.22), Inches(10.25), Inches(0.3), 16, WHITE, True, PP_ALIGN.CENTER)

    # 6 RAG detail
    slide = prs.slides.add_slide(blank); background(slide, 6, "RAG Hattı: Neden Hibrit Arama ve Reranking?")
    stages = [("PDF / MD", "Yükleme", CYAN), ("900 / 150", "Parçalama", TEAL), ("BGE", "Embedding", GREEN), ("FAISS + BM25", "Hibrit Arama", AMBER), ("Cross Encoder", "Reranking", "E95C65"), ("Qwen", "Yerel Yanıt", "A26EEA")]
    for i, (top, bottom, color) in enumerate(stages):
        x = Inches(0.55 + i * 2.1)
        add_shape(slide, Shape.ROUNDED_RECTANGLE, x, Inches(2.25), Inches(1.65), Inches(1.02), color, color, True)
        textbox(slide, top, x + Inches(0.1), Inches(2.52), Inches(1.45), Inches(0.22), 12, NAVY, True, PP_ALIGN.CENTER)
        textbox(slide, bottom, x - Inches(0.05), Inches(3.52), Inches(1.75), Inches(0.35), 11, WHITE, True, PP_ALIGN.CENTER)
        if i < len(stages) - 1:
            line(slide, x + Inches(1.67), Inches(2.76), x + Inches(2.02), Inches(2.76), MUTED, 2)
    card(slide, "Dense Search", "Anlamsal yakınlığı yakalar: kullanıcı farklı kelimeler kullansa bile uygun bağlamı bulur.", Inches(0.85), Inches(4.75), Inches(5.55), Inches(1.18), TEAL, "D")
    card(slide, "BM25 + Reranking", "Taktik terim eşleşmesini korur ve en güçlü kanıt parçalarını üst sıraya taşır.", Inches(6.9), Inches(4.75), Inches(5.55), Inches(1.18), AMBER, "R")

    # 7 security
    slide = prs.slides.add_slide(blank); background(slide, 7, "Güvenlik ve Gizlilik: On-Premise Tasarım")
    shield = add_shape(slide, Shape.HEXAGON, Inches(0.95), Inches(1.75), Inches(2.55), Inches(2.55), "102F47", TEAL)
    textbox(slide, "LOCAL\nONLY", Inches(1.23), Inches(2.37), Inches(2.0), Inches(0.72), 23, CYAN, True, PP_ALIGN.CENTER)
    rows = [("Yerel LLM", "Ollama üzerinden yanıt üretimi", GREEN), ("Yerel Model Cache", "Sorguda haricî model hub erişimi yok", CYAN), ("Koleksiyon Ayrımı", "Resmî / açık kaynak güven sınırı", AMBER), ("İzlenebilirlik", "Kaynak ve sayfa görünürlüğü", "E95C65")]
    for i, (a, b, c) in enumerate(rows):
        y = Inches(1.45 + i * 1.17)
        add_shape(slide, Shape.OVAL, Inches(4.25), y, Inches(0.53), Inches(0.53), c, c)
        textbox(slide, "✓", Inches(4.25), y + Inches(0.09), Inches(0.53), Inches(0.25), 13, NAVY, True, PP_ALIGN.CENTER)
        textbox(slide, a, Inches(4.98), y + Inches(0.02), Inches(2.3), Inches(0.25), 15, WHITE, True)
        textbox(slide, b, Inches(7.35), y + Inches(0.02), Inches(4.65), Inches(0.35), 12, MUTED)
    textbox(slide, "Üretim ortamı önerisi: RBAC • şifreli disk • audit log • imzalı indeks artefact'ları", Inches(0.9), Inches(6.22), Inches(11.5), Inches(0.32), 13, CYAN, True, PP_ALIGN.CENTER)

    # 8 demo
    slide = prs.slides.add_slide(blank); background(slide, 8, "Canlı Demo: 60 Saniyede Değer Gösterimi")
    add_shape(slide, Shape.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.55), Inches(5.75), Inches(4.75), "102A41", "102A41", True)
    textbox(slide, "1  Kapsamı seç", Inches(1.12), Inches(1.95), Inches(4.8), Inches(0.35), 17, CYAN, True)
    textbox(slide, "“Yalnızca HAVELSAN resmî kaynakları”", Inches(1.12), Inches(2.4), Inches(4.8), Inches(0.3), 13, MUTED)
    textbox(slide, "2  Soruyu sor", Inches(1.12), Inches(3.2), Inches(4.8), Inches(0.35), 17, CYAN, True)
    textbox(slide, "“ADVENT hangi taktik veri linklerini destekler?”", Inches(1.12), Inches(3.65), Inches(4.8), Inches(0.55), 13, WHITE)
    textbox(slide, "3  Kanıtı göster", Inches(1.12), Inches(4.7), Inches(4.8), Inches(0.35), 17, CYAN, True)
    textbox(slide, "Yanıt + belge adı + sayfa + reranker skoru", Inches(1.12), Inches(5.15), Inches(4.8), Inches(0.3), 13, MUTED)
    card(slide, "Beklenen Yanıt", "Link 11, Link 16, Link 22, SIMPLE, JREAP ve VMF desteği; resmî broşür kaynak gösterimi ile sunulur.", Inches(7.15), Inches(2.1), Inches(5.1), Inches(2.55), GREEN, "✓")
    textbox(slide, "Demo mesajı: Sistem yalnızca cevap vermiyor; cevabın kanıtını da sunuyor.", Inches(7.25), Inches(5.35), Inches(5.0), Inches(0.5), 15, CYAN, True, PP_ALIGN.CENTER)

    # 9 results
    slide = prs.slides.add_slide(blank); background(slide, 9, "Mevcut Durum: Doğrulanan Çıktılar")
    metrics = [("2", "Ayrı koleksiyon", CYAN), ("68", "İndeks parçası", GREEN), ("✓", "Birim testi", AMBER), ("✓", "Uçtan uca sorgu", "E95C65")]
    for i, (big, label, color) in enumerate(metrics):
        x = Inches(0.8 + i * 3.12)
        add_shape(slide, Shape.ROUNDED_RECTANGLE, x, Inches(1.85), Inches(2.45), Inches(1.85), "102F47", color, True)
        textbox(slide, big, x, Inches(2.15), Inches(2.45), Inches(0.65), 30, color, True, PP_ALIGN.CENTER)
        textbox(slide, label, x, Inches(3.05), Inches(2.45), Inches(0.3), 12, WHITE, True, PP_ALIGN.CENTER)
    card(slide, "Önemli Not", "Reranker skoru güven yüzdesi değildir; adayların sorguya göre göreli alaka sırasını gösterir.", Inches(1.55), Inches(4.75), Inches(10.25), Inches(1.0), TEAL, "i")

    # 10 roadmap
    slide = prs.slides.add_slide(blank); background(slide, 10, "Yol Haritası: Prototipten Kurumsal Ürüne")
    timeline = [("Şimdi", "Çalışan RAG", "Ayrık koleksiyonlar\nHybrid + reranking", CYAN), ("Kısa Vade", "Kalite Ölçümü", "Golden dataset\nCitation precision/recall", GREEN), ("Orta Vade", "Kurumsal Güvenlik", "RBAC • audit log\nBelge onay akışı", AMBER), ("Uzun Vade", "Sürekli İyileştirme", "Geri bildirim\nİnsan denetimli kalite", "A26EEA")]
    line(slide, Inches(1.15), Inches(3.15), Inches(12.15), Inches(3.15), MUTED, 3)
    for i, (when, title, body, color) in enumerate(timeline):
        x = Inches(0.85 + i * 3.05)
        add_shape(slide, Shape.OVAL, x + Inches(0.9), Inches(2.75), Inches(0.78), Inches(0.78), color, color)
        textbox(slide, str(i+1), x + Inches(0.9), Inches(2.94), Inches(0.78), Inches(0.25), 13, NAVY, True, PP_ALIGN.CENTER)
        textbox(slide, when, x, Inches(1.65), Inches(2.6), Inches(0.28), 12, color, True, PP_ALIGN.CENTER)
        textbox(slide, title, x, Inches(3.82), Inches(2.6), Inches(0.35), 15, WHITE, True, PP_ALIGN.CENTER)
        textbox(slide, body, x, Inches(4.35), Inches(2.6), Inches(0.55), 11.5, MUTED, False, PP_ALIGN.CENTER)

    # 11 conclusion
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(NAVY)
    add_shape(slide, Shape.RECTANGLE, 0, 0, W, Inches(0.14), TEAL)
    add_shape(slide, Shape.OVAL, Inches(9.8), Inches(1.55), Inches(2.65), Inches(2.65), "102F47", TEAL)
    textbox(slide, "✓", Inches(10.15), Inches(2.1), Inches(1.95), Inches(0.95), 48, CYAN, True, PP_ALIGN.CENTER)
    textbox(slide, "Doğru bilgiye,\ndoğru kaynakla,\ndoğru güven sınırında erişim.", Inches(0.85), Inches(1.55), Inches(7.85), Inches(1.85), 31, WHITE, True, font="Aptos Display")
    textbox(slide, "CMS-RAG Asistanı; savunma alanında bilgi erişimini hızlandırırken\ndenetlenebilirlik ve gizlilik ilkelerini mimarinin merkezinde tutar.", Inches(0.9), Inches(4.15), Inches(7.65), Inches(0.75), 16, MUTED)
    chip(slide, "KAYNAK GÖSTEREN", Inches(0.9), Inches(5.45), Inches(1.9), CYAN)
    chip(slide, "YEREL", Inches(2.98), Inches(5.45), Inches(1.1), TEAL)
    chip(slide, "GENİŞLETİLEBİLİR", Inches(4.25), Inches(5.45), Inches(1.85), GREEN)
    textbox(slide, "Teşekkürler", Inches(0.9), Inches(6.65), Inches(2), Inches(0.25), 12, CYAN, True)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT.resolve())


if __name__ == "__main__":
    build()
